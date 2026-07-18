# -*- coding: utf-8 -*-
"""Restore Slide 5's chrome (5 logos + footer placeholder) by copying them from
Slide 6, then clear the rest. Run this BEFORE patch_slide5_arch.py to rebuild the
dark cyber-machine on a slide that again has its logos/footer.

Run:  py -X utf8 restore_slide5_chrome.py  &&  py -X utf8 patch_slide5_arch.py
"""
import sys, copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"
RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (PowerPoint still holding the file?):", type(e).__name__); sys.exit(1)

src = prs.slides[5]   # slide 6 — intact chrome
dst = prs.slides[4]   # slide 5 — currently a full-bleed image

def copy_pic(shape, dst_slide):
    blip = shape._element.find(".//" + qn("a:blip"))
    rId = blip.get(qn("r:embed"))
    image_part = shape.part.related_part(rId)
    new_el = copy.deepcopy(shape._element)
    new_rId = dst_slide.part.relate_to(image_part, RT_IMAGE)
    new_el.find(".//" + qn("a:blip")).set(qn("r:embed"), new_rId)
    dst_slide.shapes._spTree.append(new_el)

# clear slide 5 (drops the full-bleed image)
for sh in list(dst.shapes):
    sh._element.getparent().remove(sh._element)

logos = 0; footer = 0
for sh in src.shapes:
    top = sh.top / 914400.0
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and top < 1.5:
        copy_pic(sh, dst); logos += 1
    elif sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and sh.has_text_frame \
            and "QUANT-A-THAN" in sh.text_frame.text:
        dst.shapes._spTree.append(copy.deepcopy(sh._element)); footer += 1

prs.save(DECK)
print("restored slide-5 chrome: %d logos, %d footer" % (logos, footer))
