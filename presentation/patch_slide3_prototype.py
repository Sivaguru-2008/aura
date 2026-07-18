# -*- coding: utf-8 -*-
"""Remake Slide 3 as the PROTOTYPE slide: the real AURA console screenshot,
annotated with numbered markers, plus a deep right-column walkthrough that maps
each console region to the live pipeline.

Keeps logos + footer; clears the old body; embeds media/console.png.
Idempotent: also removes a previously-embedded screenshot (any picture below
the header band) before re-adding.

Run:  py -X utf8 patch_slide3_prototype.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"
IMG = "../media/console.png"

INK, MUTED = "0F172A", "5B6672"
TEAL, TEAL2, NEON, NEONHI = "0E7490", "0D9488", "2DD4BF", "5EEAD4"
PALE, PALE2, BORDER = "ECF8F7", "F6FAFA", "C6E4E1"
WHITE, DARKBG, DIM = "FFFFFF", "05090E", "7FA8A5"

def _rgb(h): return RGBColor.from_string(h)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None: sp._element.remove(st)

def txt(s, x, y, w, h, content, size=8, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = [[(t, {})] for t in content.split("\n")] if isinstance(content, str) else content
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = "Arial"
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None: r._r.get_or_add_rPr().set("spc", str(sv))
    return tb

def marker(s, cx, cy, n, d=0.26):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                            Inches(d), Inches(d))
    _strip_style(sp); sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(NEON)
    sp.line.color.rgb = _rgb(DARKBG); sp.line.width = Pt(1.0)
    txt(s, cx - d / 2, cy - d / 2, d, d, str(n), size=10, color=DARKBG, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__)
    sys.exit(1)

s = prs.slides[2]  # slide 3

# 1) clear body: keep logos (pictures high on the page) + footer placeholder
for sh in list(s.shapes):
    top = sh.top / 914400.0
    if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        continue
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and top < 1.5:
        continue  # a logo
    sh._element.getparent().remove(sh._element)

# 2) kicker + title
txt(s, 0.5, 1.24, 9.0, 0.26,
    [[("02", {"color": MUTED}), ("  ·  ", {"color": MUTED}),
      ("PROTOTYPE — THE LIVE CONSOLE", {})]],
    size=10.5, color=TEAL, bold=True, spc=280)
txt(s, 0.5, 1.50, 9.2, 0.52, "Inside the console: scan → reasoned sign-off",
    size=24, color=INK, bold=True, spc=-20)

# 3) the screenshot (16:9)
PX, PY, PW = 0.5, 2.16, 6.00
PH = PW * 900.0 / 1600.0            # keep aspect -> 3.375
pic = s.shapes.add_picture(IMG, Inches(PX), Inches(PY), width=Inches(PW))
pic.line.color.rgb = _rgb(TEAL2); pic.line.width = Pt(1.25)
txt(s, PX, PY + PH + 0.06, PW, 0.22,
    [[("localhost:8000 · AURA clinical console", {"color": TEAL, "size": 7.6, "bold": True}),
      ("  —  every value on screen is live model output, offline & CPU-only",
       {"color": MUTED, "size": 7.4, "italic": True})]])

# 4) numbered markers on the console regions (fractions of the image)
regions = [  # (n, x_frac, y_frac)
    (1, 0.085, 0.46), (2, 0.31, 0.35), (3, 0.31, 0.79),
    (4, 0.62, 0.37), (5, 0.89, 0.27), (6, 0.85, 0.80)]
for n, fx, fy in regions:
    marker(s, PX + fx * PW, PY + fy * PH, n)

# 5) deep right-column walkthrough
CX, CW = 6.72, 2.78
txt(s, CX, 2.16, CW, 0.22, "HOW THE PROTOTYPE WORKS", size=9, color=TEAL,
    bold=True, spc=120)
steps = [
    (1, "WORKLIST", "14 cases triaged by uncertainty; low-confidence studies flip to ABSTAINED."),
    (2, "STUDY + SALIENCY", "Vision localizes the finding on the CXR; saliency shows where (0.99)."),
    (3, "CALIBRATED DIFFERENTIAL", "8-qubit quantum fusion, 512 shots → calibrated probabilities (99.5%)."),
    (4, "EVIDENCE → REASONING", "Each finding's contribution is explicit; hover a node for live counterfactuals."),
    (5, "SAFETY", "Epistemic · aleatoric · OOD within envelope; 90% conformal set — abstains when unsure."),
    (6, "RECOMMEND & SIGN", "EIG ranks the next best test; the doctor accepts · edits · signs — authority stays human."),
]
sy = 2.50
for n, head, body in steps:
    txt(s, CX, sy, CW, 0.50,
        [[("%d  " % n, {"color": NEONHI, "bold": True, "size": 9}),
          (head, {"color": TEAL, "bold": True, "size": 8.3, "spc": 40})],
         [(body, {"color": INK, "size": 7.3})]],
        line_spacing=1.03)
    sy += 0.505

# 6) speaker notes
s.notes_slide.notes_text_frame.text = (
    "PROTOTYPE — THE LIVE CONSOLE  ·  1:00\n"
    "This is the working product, not a mockup — every number is live output. "
    "One: the worklist is triaged by uncertainty; low-confidence studies flip to "
    "ABSTAINED instead of guessing. Two: the vision model localizes the finding "
    "and saliency shows where. Three: our 8-qubit quantum fusion, 512 shots, "
    "produces a calibrated differential — 99.5% pneumothorax here. Four: the "
    "evidence-to-reasoning graph makes every finding's contribution explicit, "
    "with live counterfactuals on hover. Five: the safety panel reports "
    "epistemic, aleatoric and OOD energy, a 90% conformal set true 90 times in "
    "100, and abstains when out of envelope. Six: the recommender ranks the next "
    "best test by information gain, and the doctor accepts, edits or signs — "
    "authority always stays human. Scan to signed report in about twenty seconds, "
    "on one screen, offline.\n\n"
    "JUDGES SHOULD FEEL: this is real, end-to-end, and safe by construction.")

prs.save(DECK)
print("saved", DECK, "— Slide 3 remade as annotated prototype slide")
