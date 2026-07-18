# -*- coding: utf-8 -*-
"""Replace Slide 5 (architecture) with the hand-drawn 4:3 render, full-bleed.

The 1440x1080 PNG matches the 10x7.5in slide exactly. Logos/footer are baked in.
Idempotent: clears all slide-5 shapes and re-adds the single full-bleed image.

Run:  py -X utf8 patch_slide5_embed.py
"""
import sys
from pptx import Presentation
from pptx.util import Inches

DECK = "AURA_Pitch_7Slides.pptx"
IMG = "slide_architecture_4x3.png"

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (PowerPoint still holding the file?):", type(e).__name__); sys.exit(1)

s = prs.slides[4]
for sh in list(s.shapes):
    sh._element.getparent().remove(sh._element)

s.shapes.add_picture(IMG, Inches(0), Inches(0),
                     width=prs.slide_width, height=prs.slide_height)

s.notes_slide.notes_text_frame.text = (
    "ARCHITECTURE & STACK  ·  0:45\n"
    "Three layered tiers: a data foundation, an application tier of independent "
    "services behind one FastAPI gateway, and the doctor-facing presentation tier. "
    "The whole point is on the board — the QPU is a module: our 8-qubit VQC docks "
    "into the FUSION bay, and a classical twin sits behind the same interface, so "
    "whichever wins the benchmark ships with zero code change. Open bays take new "
    "services the same way. Offline, no PHI leaves the site.\n\n"
    "JUDGES SHOULD FEEL: disciplined engineering — quantum is a swappable part, not a bet.")

prs.save(DECK)
print("saved", DECK, "— Slide 5 replaced with hand-drawn architecture render")
