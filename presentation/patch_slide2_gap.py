# -*- coding: utf-8 -*-
"""Add a MARKET-GAP bridge to Slide 2 (problem slide).

Shows, in the empty lower-left band: the market gap -> what AURA patches -> how.
Ties the 'no reasoning layer' root cause directly to the product.
Non-destructive (adds shapes) and idempotent (re-run replaces its own band).

Run:  py -X utf8 patch_slide2_gap.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"

INK, MUTED = "0F172A", "5B6672"
TEAL, TEAL2 = "0E7490", "0D9488"
PALE, PALE2, BORDER = "ECF8F7", "F6FAFA", "C6E4E1"
WHITE, WARN = "FFFFFF", "C05B52"

def _rgb(h): return RGBColor.from_string(h)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None: sp._element.remove(st)

def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.12):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp); sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
    return sp

def txt(s, x, y, w, h, content, size=8, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = [[(t, {})] for t in content.split("\n")] if isinstance(content, str) else content
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = "Arial"
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None: r._r.get_or_add_rPr().set("spc", str(sv))
    return tb

def arrow(s, x1, y1, x2, y2, color=TEAL2, wpt=1.6):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c); c.shadow.inherit = False
    c.line.color.rgb = _rgb(color); c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": "med", "len": "med"}))
    return c

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__)
    sys.exit(1)

s = prs.slides[1]  # slide 2

# --- idempotency: remove a previously-added gap band -------------------------
MARK = ("THE MARKET GAP", "AURA PATCHES IT", "GAP")
for sh in list(s.shapes):
    top = sh.top / 914400.0
    if sh.has_text_frame and any(m == sh.text_frame.text.strip() or m in sh.text_frame.text
                                 for m in MARK) and 4.25 < top < 5.25:
        sh._element.getparent().remove(sh._element)
    elif (sh.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE,)
          or "Connector" in type(sh).__name__) and 4.50 < top < 5.20 \
            and (sh.left / 914400.0) < 5.35:
        sh._element.getparent().remove(sh._element)

# --- draw the gap bridge -----------------------------------------------------
txt(s, 0.5, 4.34, 5.0, 0.20, "THE MARKET GAP  ·  AND HOW AURA CLOSES IT",
    size=8.5, color=TEAL, bold=True, spc=140)

# Panel A — the gap
box(s, 0.5, 4.60, 2.20, 0.58, fill=PALE2, line=WARN, lw=1.0)
txt(s, 0.66, 4.60, 1.92, 0.58,
    [[("GAP   ", {"color": WARN, "bold": True, "size": 8}),
      ("today", {"color": MUTED, "size": 6.5, "italic": True})],
     [("Every FDA-cleared tool stops at a score — the reasoning layer doesn't exist.",
       {"color": INK, "size": 7.4})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)

arrow(s, 2.74, 4.89, 3.06, 4.89, color=TEAL2, wpt=1.8)

# Panel B — what AURA patches + how
box(s, 3.10, 4.60, 2.20, 0.58, fill=PALE, line=TEAL2, lw=1.3)
txt(s, 3.26, 4.60, 1.92, 0.58,
    [[("AURA PATCHES IT", {"color": TEAL, "bold": True, "size": 8})],
     [("becomes the missing layer — ", {"color": INK, "size": 7.4}),
      ("explains · calibrates · recommends", {"color": TEAL, "size": 7.4, "bold": True})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)

prs.save(DECK)
print("saved", DECK, "— added market-gap bridge to slide 2")
