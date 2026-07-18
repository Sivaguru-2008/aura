# -*- coding: utf-8 -*-
"""Remake Slide 2 as a MARKET-GAP slide (replaces the old issue-tree/5-Whys body).

Keeps logos + footer + background; clears and rebuilds the body into three bands:
  1) the pivot   — accuracy solved -> trust is the bottleneck
  2) the gap     — decomposed into the 4 missing answers
  3) the bridge  — THE GAP (market) -> AURA PATCHES IT (how)

Idempotent: deletes every content shape (autoshape/textbox/line) and rebuilds,
so re-running yields the same slide. Pictures + footer placeholder are preserved.

Run:  py -X utf8 patch_slide2_remake.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"

INK, MUTED = "0F172A", "5B6672"
TEAL, TEAL2 = "0E7490", "0D9488"
PALE, PALE2, BORDER = "ECF8F7", "F6FAFA", "C6E4E1"
WHITE, WARN, GRAY = "FFFFFF", "C05B52", "AEB8C2"

def _rgb(h): return RGBColor.from_string(h)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None: sp._element.remove(st)

def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.10):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp); sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
    return sp

def txt(s, x, y, w, h, content, size=8, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = [[(t, {})] for t in content.split("\n")] if isinstance(content, str) else content
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = "Arial"
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None: r._r.get_or_add_rPr().set("spc", str(sv))
    return tb

def arrow(s, x1, y1, x2, y2, color=TEAL2, wpt=1.8):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c); c.shadow.inherit = False
    c.line.color.rgb = _rgb(color); c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": "med", "len": "med"}))
    return c

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__)
    sys.exit(1)

s = prs.slides[1]

# 1) clear the body: keep pictures + footer placeholder, delete everything else
KEEP = (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER)
for sh in list(s.shapes):
    if sh.shape_type not in KEEP:
        sh._element.getparent().remove(sh._element)

# 2) kicker + title
txt(s, 0.5, 1.24, 9.0, 0.26,
    [[("01", {"color": MUTED}), ("  ·  ", {"color": MUTED}),
      ("PROBLEM — THE MARKET GAP", {})]],
    size=10.5, color=TEAL, bold=True, spc=280)
txt(s, 0.5, 1.50, 9.2, 0.52, "Why hospitals reject 95%-accurate AI",
    size=26, color=INK, bold=True, spc=-20)

# 3) BAND 1 — the pivot: accuracy solved -> trust is the bottleneck
box(s, 0.5, 2.18, 4.42, 0.74, fill=PALE2, line=GRAY, lw=1.0)
txt(s, 0.70, 2.18, 4.06, 0.74,
    [[("ACCURACY   ", {"color": MUTED, "bold": True, "size": 10.5}),
      ("✓ solved since 2017", {"color": GRAY, "size": 8.5, "italic": True})],
     [("Benchmark models already match expert radiologists.",
       {"color": MUTED, "size": 8})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
box(s, 5.08, 2.18, 4.42, 0.74, fill=PALE, line=WARN, lw=1.3)
txt(s, 5.28, 2.18, 4.06, 0.74,
    [[("TRUST   ", {"color": WARN, "bold": True, "size": 10.5}),
      ("✕ the real bottleneck", {"color": WARN, "size": 8.5, "italic": True})],
     [("Without trust, even accurate models sit unused at the bedside.",
       {"color": INK, "size": 8})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# 4) BAND 2 — the gap decomposed into 4 missing answers
txt(s, 0.5, 3.08, 9.0, 0.20,
    "THE GAP, DECOMPOSED  ·  WHAT EVERY DEPLOYED TOOL LACKS",
    size=8.5, color=TEAL, bold=True, spc=140)
cards = [("No WHY", "Black-box; cognitive load + medical-liability exposure."),
         ("No HOW SURE", "No uncertainty quantification → alert fatigue."),
         ("No WHAT IF", "No counterfactual reasoning; edge cases fail silently."),
         ("No WHAT NEXT", "Static score, not a workflow-integrated pathway.")]
cw, gap = 2.13, 0.16
cx = 0.5
for title_, sub in cards:
    box(s, cx, 3.34, cw, 1.10, fill=PALE2, line=BORDER, lw=0.9)
    box(s, cx + 0.16, 3.48, 0.18, 0.05, fill=WARN, line=None, radius=0.0)  # accent tick
    txt(s, cx + 0.16, 3.58, cw - 0.30, 0.26, title_, size=10.5, bold=True, color=INK)
    txt(s, cx + 0.16, 3.86, cw - 0.28, 0.56, sub, size=7.2, color=MUTED,
        line_spacing=1.05)
    cx += cw + gap

# 5) BAND 3 — the market-gap bridge
txt(s, 0.5, 4.58, 9.0, 0.20, "THE $2.27B AI GRAVEYARD  —  AND HOW AURA ENDS IT",
    size=8.5, color=WARN, bold=True, spc=120)
box(s, 0.5, 4.82, 4.35, 0.82, fill=PALE2, line=WARN, lw=1.1)
txt(s, 0.68, 4.82, 4.02, 0.82,
    [[("950+ FDA-cleared tools — ", {"color": WARN, "bold": True, "size": 7.6}),
      ("every one a prediction calculator.", {"color": INK, "size": 7.6})],
     [("Each outputs a score and stops: no clinical reasoning layer, no trust.",
       {"color": INK, "size": 7.6})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)
arrow(s, 4.92, 5.23, 5.28, 5.23, color=TEAL2, wpt=2.0)
box(s, 5.35, 4.82, 4.15, 0.82, fill=PALE, line=TEAL2, lw=1.4)
txt(s, 5.53, 4.82, 3.82, 0.82,
    [[("AURA: THE CLINICAL REASONING ENGINE", {"color": TEAL, "bold": True, "size": 8.3})],
     [("Demystifies the read, quantifies uncertainty, drives the next action.  ",
       {"color": INK, "size": 7.6}),
      ("The doctor decides.", {"color": TEAL, "bold": True, "size": 7.6})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)

# 6) citation (now with the real market figure)
txt(s, 0.5, 6.10, 9.0, 0.32,
    "950+ FDA-cleared devices (2025) · radiology-AI market $2.27B by 2030, 24.5% "
    "CAGR (MarketsandMarkets) · India: 1 radiologist / 100k · Zech et al., PLOS Med (2018)",
    size=7.6, color=MUTED, italic=True, line_spacing=1.0)

# 7) speaker notes
s.notes_slide.notes_text_frame.text = (
    "PROBLEM — THE $2.27B AI GRAVEYARD  ·  0:50\n"
    "Accuracy is a solved problem — benchmark models already match expert "
    "radiologists. The bottleneck is trust: without it, even a 95%-accurate model "
    "sits unused at the bedside. Decompose why, and you get four failures a "
    "clinician can't accept — no WHY (black-box outputs, cognitive load, medical "
    "liability), no HOW-SURE (no uncertainty quantification, alert fatigue), no "
    "WHAT-IF (no counterfactual reasoning, silent edge-case failure), no WHAT-NEXT "
    "(a static score, not a workflow-integrated pathway). That is the $2.27B "
    "graveyard: 950+ FDA-cleared tools, every one a prediction calculator with no "
    "clinical reasoning layer. AURA is that layer — it demystifies the read, "
    "quantifies uncertainty, and drives the next action. The doctor decides.\n\n"
    "JUDGES SHOULD FEEL: accuracy is table stakes; trust is the unserved, buildable market.")

prs.save(DECK)
print("saved", DECK, "— Slide 2 remade as market-gap slide")
