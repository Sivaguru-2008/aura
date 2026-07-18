# -*- coding: utf-8 -*-
"""Build AURA Quant-a-than 2026 deck on top of the official RIT template (base.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette ----
INK    = "0F172A"   # near-black navy ink
MUTED  = "5B6672"   # gray
TEAL   = "0E7490"   # deep teal (text-safe on white)
TEAL2  = "0D9488"   # mid teal
NEON   = "2DD4BF"   # neon teal (fills, dark slides)
NEONHI = "5EEAD4"   # bright neon (dark slide accents)
PALE   = "ECF8F7"   # pale teal card fill
PALE2  = "F6FAFA"   # near-white panel
BORDER = "C6E4E1"   # card border
WHITE  = "FFFFFF"
DARKBG = "05090E"   # dark slide background
PANEL  = "0D1922"   # dark slide panel
DIM    = "7FA8A5"   # dim teal-gray text on dark
WARN   = "C05B52"   # soft red for "gap" annotations
GOLD   = "B98A2F"   # sparing warm accent

W, H = 10.0, 7.5    # slide canvas in inches (4:3)

prs = Presentation("base.pptx")
S = list(prs.slides)  # 19 slides

# ---------------------------------------------------------------- helpers ----
def _rgb(hexs):
    return RGBColor.from_string(hexs)

def _strip_style(sp):
    """Remove <p:style> — this template's theme accent leaks through it in
    PowerPoint's renderer, overriding explicit line colors."""
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)

def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.10,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None, fill_alpha=None,
        line_alpha=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha is not None:
            _set_alpha(sp.fill.fore_color._xFill, fill_alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(lw)
        ln = sp.line._get_or_add_ln()
        if dash is not None:
            d = ln.makeelement(qn("a:prstDash"), {"val": dash})
            ln.insert(list(ln).index(ln.find(qn("a:solidFill"))) + 1, d)
        if line_alpha is not None:
            sf = ln.find(qn("a:solidFill"))
            if sf is not None:
                _set_alpha(sf, line_alpha)
    return sp

def _set_alpha(fill_elm, pct):
    """pct: 0-100 opacity."""
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is None:
        return
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    clr.append(a)

def txt(s, x, y, w, h, content, size=14, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial", spc=None,
        wrap=True, line_spacing=None, space_after=None, shrink=False):
    """content: str with \n for paragraphs, OR list of paragraphs where each
    paragraph is a list of (text, opts) run tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(content, str):
        paras = [[(t, {})] for t in content.split("\n")]
    else:
        paras = [([(r, {})] if isinstance(r, str) else r) if isinstance(r, list)
                 else [(r, {})] for r in content]
        paras = []
        for item in content:
            if isinstance(item, str):
                paras.append([(item, {})])
            else:
                paras.append(item)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        for rt, opts in runs:
            r = p.add_run()
            r.text = rt
            f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.name = opts.get("font", font)
            f.color.rgb = _rgb(opts.get("color", color))
            sp_v = opts.get("spc", spc)
            if sp_v is not None:
                r._r.get_or_add_rPr().set("spc", str(sp_v))
    return tb

def kicker(s, num, tag, color=TEAL, x=0.5, y=1.24, w=9.0, align=PP_ALIGN.LEFT):
    txt(s, x, y, w, 0.26, [[(num + "  ·  ", {"color": MUTED}), (tag, {})]],
        size=10.5, color=color, bold=True, spc=280, align=align)

def title(s, text, x=0.5, y=1.50, w=9.0, size=26, color=INK, align=PP_ALIGN.LEFT):
    txt(s, x, y, w, 0.52, text, size=size, color=color, bold=True, align=align, spc=-20)

def purpose(s, text, x=0.5, y=2.05, w=9.0, align=PP_ALIGN.LEFT):
    txt(s, x, y, w, 0.30, text, size=11.5, color=MUTED, italic=True, align=align)

def cite(s, text, y=6.06):
    txt(s, 0.5, y, 9.0, 0.24, text, size=8, color=MUTED, italic=True)

def arrow(s, x1, y1, x2, y2, color=TEAL2, wpt=1.4, dash=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c)
    c.shadow.inherit = False
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    t = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(t)
    return c

def line(s, x1, y1, x2, y2, color=BORDER, wpt=1.0, dash=None, alpha=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c)
    c.shadow.inherit = False
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    if alpha is not None:
        sf = ln.find(qn("a:solidFill"))
        if sf is not None:
            _set_alpha(sf, alpha)
    return c

def glyph_circle(s, cx, cy, d, glyph, fill=PALE, ring=TEAL2, gcolor=TEAL,
                 gsize=13, bold=True, font="Arial"):
    e = box(s, cx - d / 2, cy - d / 2, d, d, fill=fill, line=ring, lw=1.0,
            shape=MSO_SHAPE.OVAL)
    tf = e.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(gsize)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = _rgb(gcolor)
    return e

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def dark_bg(s):
    box(s, -0.05, -0.05, W + 0.1, H + 0.1, fill=DARKBG, line=None, radius=0,
        shape=MSO_SHAPE.RECTANGLE)

def rings(s, cx, cy, radii, color=NEON, alphas=(30, 16, 8)):
    for r, a in zip(radii, alphas):
        e = box(s, cx - r, cy - r, 2 * r, 2 * r, fill=None, line=color, lw=1.0,
                shape=MSO_SHAPE.OVAL, line_alpha=a)

# ==============================================================================
# SLIDE 1 — title (template, already personalized). Add project name strip.
# ==============================================================================
s = S[0]
txt(s, 5.62, 4.62, 4.2, 0.62,
    [[("AURA", {"size": 21, "bold": True, "color": NEONHI, "spc": 500})],
     [("ADAPTIVE UNCERTAINTY-AWARE RADIOLOGY ARCHITECTURE",
       {"size": 8.5, "color": WHITE, "spc": 220})]],
    align=PP_ALIGN.CENTER)
notes(s, "TITLE  ·  0:15\n"
         "Good morning. I am Sivaguru, and this is AURA - the Adaptive "
         "Uncertainty-aware Radiology Architecture. Over the next ten minutes I "
         "will show you a system that treats diagnosis not as a classification "
         "problem, but as a reasoning process - with a quantum core that is "
         "mathematically earned, not marketed.\n\n"
         "JUDGES SHOULD FEEL: confident hands, zero fluff, this team knows "
         "exactly what they built.\n"
         "ANIMATION: none. Let the template breathe. Advance on the word 'earned'.")

# ==============================================================================
# SLIDE 2 — opening cinematic (dark)
# ==============================================================================
s = S[1]
dark_bg(s)
rings(s, 5.0, 3.55, [1.15, 1.85, 2.60])
txt(s, 0.5, 1.30, 9.0, 0.3, "AURA  ·  ADAPTIVE UNCERTAINTY-AWARE RADIOLOGY ARCHITECTURE",
    size=9.5, color=DIM, bold=True, spc=350, align=PP_ALIGN.CENTER)
txt(s, 0.8, 2.72, 8.4, 1.7,
    [[("The future of diagnosis is not classification.", {"color": WHITE})],
     [("It is clinical intelligence.", {"color": NEONHI, "bold": True})]],
    size=28, align=PP_ALIGN.CENTER, line_spacing=1.25)
txt(s, 0.3, 4.78, 9.4, 0.3,
    "QUANTUM-ENCODED EVIDENCE   ·   BAYESIAN REASONING AGENTS   ·   CALIBRATED UNCERTAINTY",
    size=9, color=DIM, spc=180, align=PP_ALIGN.CENTER)
line(s, 4.2, 6.42, 5.8, 6.42, color=NEON, wpt=1.0, alpha=45)
txt(s, 3.0, 6.55, 4.0, 0.26, "RIT QUANT-A-THAN 2026", size=9, color=DIM,
    spc=350, align=PP_ALIGN.CENTER)
notes(s, "OPENING  ·  0:25\n"
         "Every medical AI you have seen outputs a probability. No radiologist "
         "on earth reasons that way. A radiologist forms hypotheses, weighs "
         "evidence, quantifies doubt, and decides what to look at next. AURA is "
         "that process, built as software - with a quantum representation layer "
         "where the mathematics genuinely rewards it.\n\n"
         "JUDGES SHOULD FEEL: this is a keynote, not a slideshow. Curiosity spike.\n"
         "ANIMATION: fade transition (already set). Suggest 'It is clinical "
         "intelligence' appears on click with a slow fade (0.5 s).")

# ==============================================================================
# SLIDE 3 — the trust gap
# ==============================================================================
s = S[2]
kicker(s, "01", "THE PROBLEM")
title(s, "95% accurate. Still not trusted.")
purpose(s, "Purpose - establish that accuracy was never the bottleneck; epistemic transparency is.")

stats = [
    ("950+", "FDA-cleared AI medical devices - over 75% in radiology (2025)"),
    ("1 : 100k", "radiologist-to-population ratio in India, vs ~1 : 10k in the US"),
    ("0", "of today's deployed classifiers report when their own output should not be trusted"),
]
sy = 2.55
for big, small in stats:
    box(s, 0.5, sy, 4.05, 0.98, fill=PALE2, line=BORDER)
    txt(s, 0.72, sy + 0.17, 1.35, 0.62, big, size=25, color=TEAL, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.14, sy + 0.14, 2.28, 0.72, small, size=9.5, color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    sy += 1.14

box(s, 4.95, 2.55, 4.55, 2.66, fill=PALE, line=BORDER)
txt(s, 5.25, 2.72, 4.0, 0.26, "WHAT A CLINICIAN ACTUALLY ASKS", size=9.5,
    color=TEAL, bold=True, spc=240)
qs = ["Why this diagnosis?",
      "How certain - and which kind of uncertain?",
      "What evidence would change your mind?",
      "What should we do next?"]
qy = 3.08
for q in qs:
    glyph_circle(s, 5.42, qy + 0.135, 0.27, "?", fill=WHITE, ring=TEAL2,
                 gcolor=TEAL, gsize=10)
    txt(s, 5.68, qy, 3.7, 0.28, q, size=11.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    qy += 0.50
box(s, 4.95, 5.38, 4.55, 0.52, fill=None, line=WARN, lw=1.0, dash="dash")
txt(s, 5.15, 5.47, 4.15, 0.34,
    [[("A probability answers none of these.  ", {"color": WARN, "bold": True}),
      ("This is the trust gap.", {"color": INK, "italic": True})]], size=11)
cite(s, "Rajpurkar et al., arXiv:1711.05225 (2017)  ·  Zech et al., PLOS Medicine 15(11) (2018)  ·  US FDA, AI-Enabled Medical Device List (2025)")
notes(s, "THE PROBLEM  ·  0:50\n"
         "CheXNet passed radiologist-level F1 on pneumonia back in 2017. Eight "
         "years and nine hundred cleared devices later, most radiology AI still "
         "sits outside the clinical decision loop. Why? Because accuracy was "
         "never the bottleneck. A clinician asks four questions - why, how "
         "certain, what would change your mind, what next. A softmax score "
         "answers none of them. Zech showed these models silently collapse under "
         "distribution shift - and worse, they do not know that they collapsed. "
         "That is the trust gap AURA is built to close.\n\n"
         "JUDGES SHOULD FEEL: the problem is epistemic, not statistical - a "
         "framing they have not heard from a student team.\n"
         "ANIMATION: stats appear left-to-right on one click (wipe, 0.3 s each); "
         "red gap box last.")

# ==============================================================================
# SLIDE 4 — current pipeline & the reasoning gap
# ==============================================================================
s = S[3]
kicker(s, "02", "STATUS QUO")
title(s, "Today's pipeline ends where reasoning should begin")
purpose(s, "Purpose - locate precisely which cognitive steps current medical AI amputates.")

nodes = [("IMAGE", "chest X-ray, DICOM"), ("CNN", "f(x) - black box"),
         ("p = 0.87", "one scalar"), ("DOCTOR", "left to decide alone")]
nx, ny, nw, nh, gap = 0.55, 2.75, 1.95, 0.95, 0.37
for i, (lab, sub) in enumerate(nodes):
    x = nx + i * (nw + gap)
    fill = PALE2 if i < 3 else PALE
    b = box(s, x, ny, nw, nh, fill=fill, line=BORDER)
    txt(s, x, ny + 0.18, nw, 0.3, lab, size=14, bold=True, color=INK,
        align=PP_ALIGN.CENTER)
    txt(s, x, ny + 0.52, nw, 0.26, sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + nw + 0.04, ny + nh / 2, x + nw + gap - 0.04, ny + nh / 2)

box(s, 2.90, 4.15, 6.60, 1.42, fill=None, line=WARN, lw=1.2, dash="dash")
txt(s, 3.14, 4.28, 6.1, 0.26, "THE REASONING GAP - amputated between model and clinician",
    size=9.5, color=WARN, bold=True, spc=160)
missing = [("uncertainty decomposition", "epistemic vs aleatoric"),
           ("calibration guarantee", "is 0.87 actually 0.87?"),
           ("counterfactual justification", "what would flip the call?"),
           ("next-step planning", "which test resolves doubt?")]
mx, my = 3.14, 4.60
for i, (m, msub) in enumerate(missing):
    x = mx + (i % 2) * 3.2
    y = my + (i // 2) * 0.47
    txt(s, x, y, 3.1, 0.22,
        [[("✕  ", {"color": WARN, "bold": True, "size": 11}),
          (m, {"bold": True, "size": 10.5})]], color=INK, wrap=False)
    txt(s, x + 0.26, y + 0.225, 2.85, 0.2, msub, size=8.5, color=MUTED, wrap=False)
txt(s, 0.5, 5.72, 9.0, 0.3, "A single scalar carries no argument - and clinicians know it.",
    size=12, color=TEAL, italic=True, align=PP_ALIGN.CENTER)
cite(s, "Rudin, Nature Machine Intelligence 1, 206-215 (2019)  ·  Kompa, Snoek & Beam, npj Digital Medicine 4, 4 (2021)")
notes(s, "STATUS QUO  ·  0:40\n"
         "Here is the entire architecture of deployed medical AI: image in, CNN, "
         "one scalar out, doctor on their own. Everything in the red box - the "
         "actual cognition of diagnosis - is missing. Kompa and Beam argued in "
         "npj Digital Medicine that medicine specifically needs models that can "
         "say 'I don't know'. Rudin argued we should build interpretable "
         "reasoning rather than post-hoc excuses. AURA takes both seriously as "
         "an architecture, not as an afterthought.\n\n"
         "JUDGES SHOULD FEEL: sharp diagnosis of the field's failure mode; the "
         "red box creates an itch that slide 5 scratches.\n"
         "ANIMATION: pipeline draws left to right; red gap box fades in on "
         "second click with the four ✕ items.")

# ==============================================================================
# SLIDE 5 — the AURA pipeline (innovation)
# ==============================================================================
s = S[4]
kicker(s, "03", "OUR INNOVATION")
title(s, "AURA: diagnosis as a reasoning pipeline")
purpose(s, "Purpose - replace the classifier with an eight-stage diagnostic operating system.")

stages_r1 = [("01", "Image", "DICOM study + context"),
             ("02", "Quantum encoding", "feature map Uφ(x)"),
             ("03", "Bayesian agents", "belief graph over DDx"),
             ("04", "Calibration", "conformal + temperature")]
stages_r2 = [("08", "Clinician", "decides - with evidence"),
             ("07", "Grounded report", "every claim cited"),
             ("06", "Evidence planner", "next test by EIG"),
             ("05", "Counterfactuals", "what flips the call")]
nw, nh, gap = 2.02, 1.06, 0.28
r1y, r2y = 2.62, 4.42
for i, (num, lab, sub) in enumerate(stages_r1):
    x = 0.55 + i * (nw + gap)
    hi = (i == 1)
    b = box(s, x, r1y, nw, nh, fill=(PALE if hi else PALE2),
            line=(TEAL2 if hi else BORDER), lw=1.2 if hi else 0.75)
    txt(s, x + 0.14, r1y + 0.10, 0.5, 0.22, num, size=9, color=TEAL, bold=True, spc=100)
    txt(s, x + 0.14, r1y + 0.32, nw - 0.28, 0.30, lab, size=12.5, bold=True, color=INK)
    txt(s, x + 0.14, r1y + 0.64, nw - 0.28, 0.34, sub, size=9, color=MUTED)
    if i < 3:
        arrow(s, x + nw + 0.03, r1y + nh / 2, x + nw + gap - 0.03, r1y + nh / 2)
for i, (num, lab, sub) in enumerate(stages_r2):
    x = 0.55 + i * (nw + gap)
    hi = (i == 0)
    b = box(s, x, r2y, nw, nh, fill=(PALE if hi else PALE2),
            line=(TEAL2 if hi else BORDER), lw=1.2 if hi else 0.75)
    txt(s, x + 0.14, r2y + 0.10, 0.5, 0.22, num, size=9, color=TEAL, bold=True, spc=100)
    txt(s, x + 0.14, r2y + 0.32, nw - 0.28, 0.30, lab, size=12.5, bold=True, color=INK)
    txt(s, x + 0.14, r2y + 0.64, nw - 0.28, 0.34, sub, size=9, color=MUTED)
    if 0 < i < 3:
        arrow(s, x + nw + gap - 0.03, r2y + nh / 2, x + nw + 0.03, r2y + nh / 2)
# connect row1 end -> row2 end (down the right side)
lastx = 0.55 + 3 * (nw + gap)
arrow(s, lastx + nw / 2, r1y + nh + 0.04, lastx + nw / 2, r2y - 0.04)
txt(s, 0.5, 5.72, 9.0, 0.3, "Not a better classifier - a diagnostic operating system.",
    size=12.5, color=TEAL, italic=True, bold=False, align=PP_ALIGN.CENTER)
cite(s, "Pipeline implemented end-to-end: services/{vision, fusion, safety, explain, recommend, report} - live demo available.")
notes(s, "OUR INNOVATION  ·  1:00\n"
         "This is AURA. The image enters, but it is never handed to a lone "
         "classifier. Stage two embeds compressed clinical evidence into a "
         "quantum feature space. Stage three: reasoning agents run Bayesian "
         "updates over a differential-diagnosis belief graph. Stage four "
         "calibrates - conformal sets with a coverage guarantee. Then the system "
         "asks its own follow-ups: what would change my mind, and which test "
         "buys the most information for this patient. The output is a grounded "
         "report where every sentence is linked to evidence. The clinician stays "
         "the decision-maker - stage eight, not a rubber stamp.\n"
         "Point at the boxes as you walk the flow. Do NOT read the boxes.\n\n"
         "JUDGES SHOULD FEEL: architectural ambition with discipline; each "
         "stage is a real subsystem they will see again.\n"
         "ANIMATION: row 1 wipes in, then the down-arrow, then row 2 "
         "right-to-left, mirroring the snake flow.")

# ==============================================================================
# SLIDE 6 — system architecture
# ==============================================================================
s = S[5]
kicker(s, "04", "SYSTEM ARCHITECTURE")
title(s, "Five layers, nine replaceable services")
purpose(s, "Purpose - show this is an engineered system with typed contracts, not a notebook.")

layers = [
    ("CLINICAL", "Doctor dashboard · differential view · report signing · feedback capture", False),
    ("KNOWLEDGE", "Disease ontology · guideline engine · case memory · audit ledger", False),
    ("AI REASONING", "Vision encoder · Bayesian reasoning agents · uncertainty + safety engine", False),
    ("QUANTUM", "Feature-map registry · quantum kernels · VQC evidence fusion · quantum similarity", True),
    ("GATEWAY / DATA", "FastAPI gateway · async event bus · PACS-EHR adapters · SQLite provenance", False),
]
ly = 2.50
for name, desc, hi in layers:
    box(s, 0.5, ly, 6.95, 0.60, fill=(PALE if hi else PALE2),
        line=(TEAL2 if hi else BORDER), lw=1.3 if hi else 0.75)
    txt(s, 0.72, ly + 0.06, 1.75, 0.48, name, size=10.5, bold=True,
        color=(TEAL if hi else INK), spc=140, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.55, ly + 0.06, 4.75, 0.48, desc, size=9.5, color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE)
    ly += 0.68
rails = [("SAFETY", "abstention policy - no silent failure"),
         ("PROVENANCE", "every claim carries {model, seed, version}"),
         ("CONTRACTS", "shared Pydantic schemas across services")]
ry = 2.50
for name, desc in rails:
    box(s, 7.62, ry, 1.88, 1.04, fill=WHITE, line=BORDER)
    txt(s, 7.76, ry + 0.12, 1.6, 0.22, name, size=8.5, bold=True, color=TEAL, spc=180)
    txt(s, 7.76, ry + 0.36, 1.62, 0.62, desc, size=8.5, color=MUTED, line_spacing=1.0)
    ry += 1.16
txt(s, 0.5, 5.98, 7.0, 0.26,
    "Every engine is independently swappable - the quantum layer ships beside its classical twin.",
    size=10, color=TEAL, italic=True)
notes(s, "ARCHITECTURE  ·  0:40\n"
         "AURA is nine microservices over an async event bus, speaking shared "
         "typed contracts. The quantum layer is one layer - deliberately not the "
         "whole story. And note the right rail: safety, provenance, contracts "
         "are cross-cutting. Every posterior AURA has ever produced is "
         "reproducible from its provenance record - backend, seed, circuit "
         "version, shots. That is what auditability means in a hospital.\n\n"
         "JUDGES SHOULD FEEL: production engineering maturity - startup founders "
         "will notice the swappable-services decision.\n"
         "ANIMATION: layers stack bottom-up (gateway first), quantum layer "
         "pulses last; rail fades in together.")

# ==============================================================================
# SLIDE 7 — quantum core
# ==============================================================================
s = S[6]
kicker(s, "05", "QUANTUM CORE")
title(s, "Why Hilbert space - the mathematical case")
purpose(s, "Purpose - justify the quantum layer with geometry, never with speed.")

steps = [
    ("x ∈ R⁸", "compressed evidence vector - 8 clinical channels, never raw pixels"),
    ("Uφ(x) |0⟩⊗⁸", "parameterized quantum circuit angle-encodes evidence"),
    ("|φ(x)⟩ ∈ H,  dim = 2⁸", "entangling layers spread data across a 256-dim state space"),
    ("k(x, x′) = |⟨φ(x)|φ(x′)⟩|²", "quantum kernel - fidelity between diagnostic states"),
    ("p(d | e)  via VQC fusion", "hybrid Bayesian update over the differential"),
]
syy = 2.52
for i, (eq, desc) in enumerate(steps):
    box(s, 0.5, syy, 3.95, 0.56, fill=PALE2, line=BORDER)
    txt(s, 0.68, syy + 0.055, 3.6, 0.24, eq, size=11.5, bold=True, color=TEAL,
        font="Cambria")
    txt(s, 0.68, syy + 0.30, 3.62, 0.22, desc, size=7.8, color=MUTED)
    if i < 4:
        arrow(s, 2.47, syy + 0.56, 2.47, syy + 0.66, wpt=1.2)
    syy += 0.665

box(s, 4.75, 2.52, 4.75, 1.55, fill=PALE, line=TEAL2, lw=1.1)
txt(s, 5.0, 2.66, 4.3, 0.26, "ENTANGLING FEATURE MAP", size=9, color=TEAL,
    bold=True, spc=220)
txt(s, 5.0, 2.94, 4.3, 0.42,
    "Uφ(x) = exp( i Σ_S φ_S(x) Π_k∈S Z_k )",
    size=15, color=INK, bold=True, font="Cambria", align=PP_ALIGN.CENTER)
txt(s, 5.0, 3.44, 4.28, 0.56,
    "Second-order terms φ_{jk}(x) = (π - x_j)(π - x_k) encode pairwise "
    "evidence interactions directly in the state - e.g. opacity × fever.",
    size=9, color=MUTED, line_spacing=1.05)

box(s, 4.75, 4.22, 4.75, 1.60, fill=WHITE, line=BORDER)
txt(s, 5.0, 4.34, 4.3, 0.26, "THE HONEST CLAIM", size=9, color=TEAL, bold=True, spc=220)
txt(s, 5.0, 4.62, 4.28, 1.12,
    [[("Richer geometry, not speed. ", {"bold": True, "color": INK}),
      ("Kernels from entangling circuits are conjectured classically hard to "
       "estimate (Havlíček 2019); a rigorous separation exists for a "
       "related family (Liu 2021). We claim the ", {"color": MUTED})],
     [("representation", {"italic": True, "color": INK}),
      (" - and benchmark it against a classical twin on every metric.", {"color": MUTED})]],
    size=9.5, line_spacing=1.1)
cite(s, "Havlíček et al., Nature 567, 209 (2019)  ·  Schuld & Killoran, PRL 122, 040504 (2019)  ·  Liu, Arunachalam & Temme, Nature Physics 17, 1013 (2021)")
notes(s, "QUANTUM CORE  ·  1:30  - the slide the quantum researchers will grill\n"
         "We never put pixels in a circuit. Vision stays classical - CNNs won "
         "imaging. What goes into the quantum layer is a compressed 8-channel "
         "evidence vector. Eight qubits, angle encoding, entangling layers: the "
         "state lives in a 256-dimensional Hilbert space, and the ZZ-style map "
         "writes pairwise evidence interactions into the phase structure. The "
         "kernel is state fidelity - similarity in that space. Havlicek showed "
         "such kernels are conjectured hard to estimate classically; Liu, "
         "Arunachalam and Temme proved a rigorous separation for a related "
         "problem class. So our claim is representational: a richer similarity "
         "geometry over clinical evidence. Not 'quantum is faster'. And because "
         "conjecture is not evidence, every quantum service ships beside a "
         "classical twin and must beat it head-to-head - next slides show how.\n\n"
         "JUDGES SHOULD FEEL: rare intellectual honesty + real command of the "
         "QML literature. This slide wins the quantum jury.\n"
         "ANIMATION: left pipeline draws downward step by step; right cards on "
         "separate clicks. Practice this slide the most.")

# ==============================================================================
# SLIDE 8 — clinical reasoning graph
# ==============================================================================
s = S[7]
kicker(s, "06", "REASONING ENGINE")
title(s, "A belief graph, not a softmax")
purpose(s, "Purpose - show diagnosis as message passing between evidence and hypotheses.")

# evidence nodes (left)
ev = [("opacity gradient", "vision"), ("costophrenic blunting", "vision"),
      ("post-op context", "prior"), ("cohort match 0.91", "quantum similarity")]
ey = 2.62
for lab, src in ev:
    box(s, 0.5, ey, 2.30, 0.62, fill=WHITE, line=BORDER)
    txt(s, 0.64, ey + 0.08, 2.05, 0.24, lab, size=9.5, bold=True, color=INK)
    txt(s, 0.64, ey + 0.33, 2.05, 0.2, src.upper(), size=7, color=TEAL, spc=160)
    ey += 0.80

# center belief node
box(s, 3.85, 3.55, 1.85, 1.85, fill=PALE, line=TEAL2, lw=1.4, shape=MSO_SHAPE.OVAL)
txt(s, 3.85, 4.08, 1.85, 0.42,
    [[("BELIEF STATE", {"size": 9, "bold": True, "color": TEAL, "spc": 160})],
     [("p(d | e)", {"size": 14, "bold": True, "color": INK, "font": "Cambria"})]],
    align=PP_ALIGN.CENTER)
for i in range(4):
    yy = 2.62 + i * 0.80 + 0.31
    arrow(s, 2.82, yy, 3.95, 4.15 + (i - 1.5) * 0.28, color=TEAL2, wpt=1.0)

# disease posterior bars (right)
dd = [("Pleural effusion", 0.58), ("Nodule", 0.21), ("Pneumothorax", 0.12),
      ("COPD pattern", 0.09)]
dy = 2.62
for lab, p in dd:
    box(s, 6.35, dy, 3.15, 0.62, fill=PALE2, line=BORDER)
    txt(s, 6.50, dy + 0.07, 2.2, 0.24, lab, size=10, bold=True, color=INK)
    txt(s, 8.62, dy + 0.07, 0.8, 0.24, f"{p:.2f}", size=10, bold=True,
        color=TEAL, align=PP_ALIGN.RIGHT)
    box(s, 6.50, dy + 0.38, 2.9 * p / 0.6, 0.09, fill=NEON, line=None, radius=0.5)
    arrow(s, 5.72, 4.30 + (dy - 4.0) * 0.0, 6.30, dy + 0.31, color=TEAL2, wpt=1.0) if False else None
    arrow(s, 5.66, 4.44 + ({0.58: -0.9, 0.21: -0.3, 0.12: 0.3, 0.09: 0.9}[p]) * 0.28,
          6.28, dy + 0.31, color=TEAL2, wpt=1.0)
    dy += 0.80

box(s, 0.5, 5.62, 9.0, 0.44, fill=None, line=None)
txt(s, 0.5, 5.62, 9.0, 0.44,
    [[("Bayesian update per message:   ", {"color": MUTED, "size": 10})],
     ],
    size=10)
txt(s, 2.62, 5.56, 4.8, 0.4,
    "p(d | e₁..eₖ) ∝ p(d) Πₖ p(eₖ | d)",
    size=14, color=TEAL, bold=True, font="Cambria", align=PP_ALIGN.CENTER)
txt(s, 0.5, 5.60, 2.1, 0.3, "Update rule", size=9, color=MUTED, italic=True,
    align=PP_ALIGN.RIGHT) if False else None
cite(s, "Koller & Friedman, Probabilistic Graphical Models, MIT Press (2009)  ·  Richens et al., Nature Communications 11, 3923 (2020) - counterfactual diagnosis")
notes(s, "REASONING ENGINE  ·  0:45\n"
         "Instead of one softmax, AURA maintains a belief graph. Evidence nodes "
         "arrive from three sources - the vision encoder, patient priors, and "
         "quantum similarity against precedent cases. Reasoning agents pass "
         "messages: each piece of evidence multiplies into the posterior over "
         "the differential. The output is not 'effusion 0.58' - it is a ranked "
         "differential with the evidence trail that produced it. Richens showed "
         "at DeepMind that causal, counterfactual diagnosis beats associative "
         "diagnosis - this graph is what makes that possible here.\n\n"
         "JUDGES SHOULD FEEL: AI researchers recognize PGM rigor; clinicians "
         "recognize their own differential process.\n"
         "ANIMATION: evidence nodes in first, edges draw toward belief state, "
         "then posterior bars grow (0.4 s bar-grow).")

# ==============================================================================
# SLIDE 9 — explainability
# ==============================================================================
s = S[8]
kicker(s, "07", "EXPLAINABILITY")
title(s, "Four artifacts behind every claim")
purpose(s, "Purpose - explanation as evidence, not decoration; each card is a live output.")

cards = [
    ("WHERE", "Occlusion saliency", "model-agnostic perturbation map - which regions drive the finding", "saliency"),
    ("WHAT IF", "Counterfactual regions",
     "δ* = argmin ‖δ‖  s.t.  diagnosis(x + δ) ≠ diagnosis(x)", "cf"),
    ("WHY", "Quantum precedent search", "nearest prior cases by state fidelity  Tr[ρ(x)ρ(xᵢ)]", "sim"),
    ("HOW SURE", "Conformal intervals", "sets, not points - {effusion, atelectasis} at 90% coverage", "conf"),
]
positions = [(0.5, 2.52), (5.05, 2.52), (0.5, 4.28), (5.05, 4.28)]
for (tag, name, desc, kind), (x, y) in zip(cards, positions):
    box(s, x, y, 4.45, 1.62, fill=PALE2, line=BORDER)
    txt(s, x + 0.22, y + 0.14, 1.6, 0.22, tag, size=8.5, color=TEAL, bold=True, spc=240)
    txt(s, x + 0.22, y + 0.38, 2.9, 0.28, name, size=13, bold=True, color=INK)
    txt(s, x + 0.22, y + 0.72, 2.85, 0.74, desc, size=9,
        color=MUTED, line_spacing=1.1,
        font=("Cambria" if kind in ("cf", "sim") else "Arial"))
    mx = x + 3.22
    if kind == "saliency":
        box(s, mx, y + 0.28, 1.02, 1.08, fill=INK, line=None, radius=0.08)
        box(s, mx + 0.18, y + 0.52, 0.5, 0.42, fill=NEON, line=None,
            shape=MSO_SHAPE.OVAL, fill_alpha=55)
        box(s, mx + 0.30, y + 0.62, 0.26, 0.22, fill=NEONHI, line=None,
            shape=MSO_SHAPE.OVAL)
    elif kind == "cf":
        box(s, mx, y + 0.28, 1.02, 1.08, fill=INK, line=None, radius=0.08)
        box(s, mx + 0.22, y + 0.50, 0.42, 0.38, fill=None, line=WARN, lw=1.2,
            dash="dash", shape=MSO_SHAPE.OVAL)
        txt(s, mx, y + 0.98, 1.02, 0.3, "remove → flips", size=6.5,
            color=WHITE, align=PP_ALIGN.CENTER)
    elif kind == "sim":
        for k in range(3):
            box(s, mx + 0.06 + k * 0.30, y + 0.36 + k * 0.16, 0.44, 0.44,
                fill=PALE, line=TEAL2, lw=0.75, shape=MSO_SHAPE.OVAL)
        txt(s, mx - 0.06, y + 1.06, 1.2, 0.3, "0.93 · 0.91 · 0.88", size=7.5,
            color=TEAL, align=PP_ALIGN.CENTER, bold=True)
    elif kind == "conf":
        for k, wfrac in enumerate([0.85, 0.55, 0.3]):
            box(s, mx, y + 0.40 + k * 0.30, 1.0 * wfrac, 0.12, fill=NEON,
                line=None, radius=0.5, fill_alpha=90 - k * 25)
txt(s, 0.5, 6.0, 9.0, 0.26,
    "If a claim has no artifact, AURA does not make the claim.",
    size=10.5, color=TEAL, italic=True, align=PP_ALIGN.CENTER) if False else None
cite(s, "Selvaraju et al., ICCV (2017)  ·  Wachter et al., Harvard JOLT 31 (2018)  ·  Zeiler & Fergus, ECCV (2014)")
notes(s, "EXPLAINABILITY  ·  0:45\n"
         "Four artifacts ship with every report. Where - perturbation saliency, "
         "chosen over pure gradient methods because it is model-agnostic and "
         "auditable. What-if - the minimal image region whose removal flips the "
         "differential; that is a falsifiable explanation, in Wachter's "
         "counterfactual framework. Why - precedent: the nearest prior cases in "
         "the quantum latent space, by state fidelity - doctors reason by "
         "precedent, so does AURA. And how sure - intervals, never points. Rule "
         "of the house: if a claim has no artifact, the report does not say it.\n\n"
         "JUDGES SHOULD FEEL: explainability with falsifiability - beyond "
         "heat-map theater.\n"
         "ANIMATION: cards appear clockwise from WHERE; keep 0.3 s fades.")

# ==============================================================================
# SLIDE 10 — uncertainty
# ==============================================================================
s = S[9]
kicker(s, "08", "UNCERTAINTY ENGINE")
title(s, "Knowing which kind of 'unsure'")
purpose(s, "Purpose - decompose doubt, because each kind demands a different clinical action.")

box(s, 0.5, 2.48, 9.0, 0.66, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 0.5, 2.58, 9.0, 0.44,
    [[("H[ p(y|x) ]", {"color": INK, "bold": True}),
      ("   =   ", {"color": MUTED}),
      ("I(y ; θ | x)", {"color": TEAL, "bold": True}),
      ("   +   ", {"color": MUTED}),
      ("Eθ H[ p(y|x,θ) ]", {"color": TEAL2, "bold": True}),
      ("        total = epistemic + aleatoric", {"color": MUTED, "size": 10, "italic": True})]],
    size=15, font="Cambria", align=PP_ALIGN.CENTER)

ucards = [
    ("EPISTEMIC", "model doubt", "deep-ensemble + MC-dropout disagreement; shrinks with data",
     "→  defer to radiologist"),
    ("ALEATORIC", "data noise", "irreducible ambiguity in the image itself",
     "→  recommend re-imaging"),
    ("OUT-OF-DIST.", "wrong input", "energy-score OOD + quantum latent density",
     "→  abstain, flag input"),
]
ux = 0.5
for name, tag, desc, act in ucards:
    box(s, ux, 3.36, 2.92, 1.52, fill=PALE2, line=BORDER)
    txt(s, ux + 0.18, 3.50, 2.6, 0.24,
        [[(name + "  ", {"bold": True, "color": INK, "size": 11}),
          ("- " + tag, {"color": MUTED, "size": 9, "italic": True})]])
    txt(s, ux + 0.18, 3.80, 2.58, 0.55, desc, size=8.8, color=MUTED, line_spacing=1.05)
    txt(s, ux + 0.18, 4.48, 2.6, 0.26, act, size=9.5, color=TEAL, bold=True)
    ux += 3.04

box(s, 0.5, 5.06, 4.75, 0.92, fill=WHITE, line=BORDER)
txt(s, 0.68, 5.16, 4.4, 0.26, "CONFORMAL GUARANTEE - distribution-free", size=8.5,
    color=TEAL, bold=True, spc=180)
txt(s, 0.68, 5.42, 4.4, 0.5,
    [[("C(x) = { y : s(x,y) ≤ q̂₁₋α }", {"font": "Cambria", "bold": True, "color": INK, "size": 12.5}),
      ("     P( y ∈ C(x) ) ≥ 1 − α", {"font": "Cambria", "color": TEAL, "size": 12.5, "bold": True})]],
    size=12.5)
# mini calibration curve
box(s, 5.5, 5.06, 4.0, 0.92, fill=WHITE, line=BORDER)
txt(s, 5.68, 5.16, 3.6, 0.24, "CALIBRATION - temperature-scaled", size=8.5,
    color=TEAL, bold=True, spc=180)
line(s, 5.85, 5.86, 7.05, 5.86, color=MUTED, wpt=0.75)     # x axis
line(s, 5.85, 5.86, 5.85, 5.42, color=MUTED, wpt=0.75)     # y axis
line(s, 5.85, 5.86, 7.05, 5.44, color=BORDER, wpt=1.0, dash="dash")  # ideal
pts_curve = [(5.85, 5.86), (6.15, 5.79), (6.45, 5.68), (6.75, 5.55), (7.05, 5.46)]
for (cx1, cy1), (cx2, cy2) in zip(pts_curve, pts_curve[1:]):
    line(s, cx1, cy1, cx2, cy2, color=NEON, wpt=1.75)
txt(s, 7.25, 5.44, 2.1, 0.5,
    [[("ECE  0.142 → 0.031", {"bold": True, "color": INK, "size": 10.5})],
     [("after temperature scaling", {"color": MUTED, "size": 8})]])
cite(s, "Kendall & Gal, NeurIPS (2017)  ·  Guo et al., ICML (2017)  ·  Angelopoulos & Bates, Foundations and Trends in ML 16(4) (2023)")
notes(s, "UNCERTAINTY  ·  1:00\n"
         "One equation runs this slide: total predictive entropy decomposes into "
         "epistemic - mutual information between prediction and parameters - "
         "plus aleatoric. Why decompose? Because each kind of doubt maps to a "
         "different clinical action. Epistemic high: the model is under-trained "
         "for this case - defer to the human. Aleatoric high: the image itself "
         "is ambiguous - more data will not help, re-image. Out-of-distribution: "
         "abstain entirely. And on top, conformal prediction gives a "
         "distribution-free coverage guarantee - the true label is in the "
         "predicted set with probability at least 1 minus alpha, by theorem, "
         "not by hope. ECE numbers are from our own benchmark harness.\n\n"
         "JUDGES SHOULD FEEL: statistical literacy at research level - the "
         "uncertainty-to-action mapping is the memorable move.\n"
         "ANIMATION: equation first, three cards left-to-right, guarantee last.")

# ==============================================================================
# SLIDE 11 — next best evidence
# ==============================================================================
s = S[10]
kicker(s, "09", "EVIDENCE PLANNER")
title(s, "The next question, not just the current answer")
purpose(s, "Purpose - turn residual uncertainty into an optimal, cost-aware next test.")

box(s, 0.5, 2.50, 4.1, 1.55, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 0.72, 2.64, 3.7, 0.26, "EXPECTED INFORMATION GAIN", size=9, color=TEAL,
    bold=True, spc=200)
txt(s, 0.72, 2.94, 3.7, 0.40, "a* = argmaxₐ  EIG(a) − λ·cost(a)",
    size=14.5, color=INK, bold=True, font="Cambria")
txt(s, 0.72, 3.40, 3.7, 0.36, "EIG(a) = H[p(d)] − Eₒ H[p(d | o)]",
    size=12.5, color=TEAL2, font="Cambria", bold=True)
txt(s, 0.72, 3.78, 3.68, 0.24, "computed through the fusion posterior, per patient",
    size=8.5, color=MUTED, italic=True)

tests = [("CT thorax", 0.82, "high cost"), ("Repeat PA X-ray", 0.64, "low cost"),
         ("Thoracic ultrasound", 0.51, "low cost"), ("CBC + CRP panel", 0.33, "minimal"),
         ("Observe 24 h", 0.12, "no cost")]
txt(s, 4.95, 2.50, 4.55, 0.26, "RANKED FOR THIS PATIENT  (bits of entropy resolved)",
    size=9, color=TEAL, bold=True, spc=160)
ty = 2.86
for name, eig, cost in tests:
    txt(s, 4.95, ty, 1.85, 0.24, name, size=10, bold=True, color=INK,
        anchor=MSO_ANCHOR.MIDDLE)
    box(s, 6.85, ty + 0.015, 1.9 * eig / 0.82, 0.19, fill=NEON, line=None,
        radius=0.5, fill_alpha=100 - int((0.82 - eig) * 60))
    txt(s, 8.85, ty, 0.62, 0.24, f"{eig:.2f}", size=10, bold=True, color=TEAL,
        anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.42
box(s, 4.95, ty + 0.06, 4.55, 0.52, fill=PALE2, line=BORDER)
txt(s, 5.12, ty + 0.135, 4.2, 0.32,
    [[("Recommendation:  ", {"color": MUTED, "size": 9.5}),
      ("repeat PA X-ray ", {"bold": True, "color": INK, "size": 10.5}),
      ("- 78% of CT's information at 4% of its cost", {"color": TEAL, "size": 9.5, "italic": True})]])
txt(s, 0.5, 5.62, 9.0, 0.3,
    "AURA optimizes the diagnostic trajectory - Lindley's 1956 experimental design, made clinical.",
    size=11, color=TEAL, italic=True, align=PP_ALIGN.CENTER)
cite(s, "Lindley, Ann. Math. Statist. 27(4), 986-1005 (1956)  ·  Bernardo, Ann. Statist. 7(3) (1979)  ·  implemented in services/recommend")
notes(s, "EVIDENCE PLANNER  ·  0:50\n"
         "Here is where AURA stops being a classifier at all. Given the current "
         "posterior, it evaluates candidate actions - CT, repeat X-ray, "
         "ultrasound, labs, observation - by expected information gain: the "
         "expected entropy reduction of the differential, penalized by cost and "
         "risk. In this case: CT resolves the most entropy, but a repeat PA "
         "film gets 78 percent of the information at 4 percent of the cost - so "
         "that is the recommendation. This is Lindley's Bayesian experimental "
         "design from 1956, running per-patient. Fewer unnecessary scans, "
         "faster convergence to the answer.\n\n"
         "JUDGES SHOULD FEEL: 'the AI asks for the right test' - startup "
         "founders see the health-economics value instantly.\n"
         "ANIMATION: equation first; bars grow staggered; recommendation chip "
         "pops last.")

# ==============================================================================
# SLIDE 12 — live workflow loop
# ==============================================================================
s = S[11]
kicker(s, "10", "CLINICAL LOOP")
title(s, "Human-in-the-loop by construction")
purpose(s, "Purpose - place AURA inside the hospital workflow, learning from every sign-off.")

import math
cx, cy, R = 4.05, 4.12, 1.42
loop = [("PATIENT", "presents"), ("X-RAY / PACS", "study acquired"),
        ("AURA", "report + uncertainty"), ("RADIOLOGIST", "reviews, signs"),
        ("EHR / WARD", "decision executed"), ("FEEDBACK", "accept / correct")]
pts = []
for i in range(6):
    ang = -90 + i * 60
    px = cx + R * math.cos(math.radians(ang))
    py = cy + R * math.sin(math.radians(ang))
    pts.append((px, py))
for i, ((px, py), (lab, sub)) in enumerate(zip(pts, loop)):
    hi = (lab == "AURA")
    bw, bh = 1.42, 0.60
    box(s, px - bw / 2, py - bh / 2, bw, bh, fill=(PALE if hi else WHITE),
        line=(TEAL2 if hi else BORDER), lw=1.3 if hi else 0.75)
    txt(s, px - bw / 2, py - 0.225, bw, 0.24, lab, size=9.5, bold=True,
        color=(TEAL if hi else INK), align=PP_ALIGN.CENTER)
    txt(s, px - bw / 2, py + 0.015, bw, 0.2, sub, size=7.5, color=MUTED,
        align=PP_ALIGN.CENTER)
for i in range(6):
    x1, y1 = pts[i]
    x2, y2 = pts[(i + 1) % 6]
    # shrink toward midpoints so arrows do not overlap boxes
    fx = x1 + (x2 - x1) * 0.36
    fy = y1 + (y2 - y1) * 0.36
    tx = x1 + (x2 - x1) * 0.64
    ty2 = y1 + (y2 - y1) * 0.64
    arrow(s, fx, fy, tx, ty2, color=TEAL2, wpt=1.3)
glyph_circle(s, cx, cy, 1.08, "", fill=PALE2, ring=BORDER)
txt(s, cx - 0.85, cy - 0.30, 1.7, 0.6,
    [[("LEARNING LOOP", {"size": 8, "bold": True, "color": TEAL, "spc": 140})],
     [("priors · calibration", {"size": 7.5, "color": MUTED})],
     [("case memory", {"size": 7.5, "color": MUTED})]],
    align=PP_ALIGN.CENTER)

chips = [("Radiologist signs every report", "authority never delegated"),
         ("Every claim in the audit ledger", "reproducible by construction"),
         ("Feedback re-fits calibration nightly", "the system stays honest")]
cyy = 2.9
for c1, c2 in chips:
    box(s, 6.55, cyy, 2.95, 0.78, fill=PALE2, line=BORDER)
    txt(s, 6.72, cyy + 0.10, 2.62, 0.28, c1, size=9.5, bold=True, color=INK,
        line_spacing=1.0)
    txt(s, 6.72, cyy + 0.42, 2.62, 0.24, c2, size=8, color=MUTED, italic=True)
    cyy += 0.92
cite(s, "Deployment pattern follows Sendak et al., npj Digital Medicine 3, 41 (2020) - clinical ML delivery.")
notes(s, "CLINICAL LOOP  ·  0:35\n"
         "AURA sits inside the workflow, not on top of it. Study lands from "
         "PACS, AURA drafts the grounded report with its uncertainty and next "
         "test, the radiologist reviews and signs - authority is never "
         "delegated. The sign-off itself is training signal: accepts and "
         "corrections re-fit the calibration and update case memory nightly. "
         "The system gets more honest with every use.\n\n"
         "JUDGES SHOULD FEEL: deployable, regulator-aware, safe - clinicians "
         "keep the pen.\n"
         "ANIMATION: loop arrows animate clockwise once (0.2 s each) - a single "
         "rotation, then stop.")

# ==============================================================================
# SLIDE 13 — research novelty table
# ==============================================================================
s = S[12]
kicker(s, "11", "POSITIONING")
title(s, "What exists, and what didn't - until now")
purpose(s, "Purpose - locate AURA's contribution against the strongest published baselines.")

cols = ["Capability", "CheXNet\n(2017)", "MedPaLM-M\n(2023)", "MONAI\n(2020)", "BioGPT\n(2022)", "AURA"]
rows = [
    ("Disease probability", "1", "1", "1", "1", "1"),
    ("Uncertainty decomposition", "0", "0", "~", "0", "1"),
    ("Coverage-guaranteed sets", "0", "0", "0", "0", "1"),
    ("Counterfactual explanation", "0", "0", "0", "0", "1"),
    ("Next-test planning (EIG)", "0", "~", "0", "0", "1"),
    ("Quantum-kernel latent space", "0", "0", "0", "0", "1"),
    ("Grounded, auditable report", "0", "~", "0", "~", "1"),
]
tx0, ty0 = 0.5, 2.44
cw = [2.85, 1.22, 1.22, 1.22, 1.22, 1.27]
rh = 0.40
xacc = [tx0]
for wdt in cw:
    xacc.append(xacc[-1] + wdt)
# header
for j, htx in enumerate(cols):
    hi = (j == 5)
    box(s, xacc[j], ty0, cw[j], rh + 0.06, fill=(TEAL if hi else INK), line=None,
        radius=0.06)
    txt(s, xacc[j], ty0 + 0.05, cw[j], rh - 0.02, htx, size=8.5, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER, line_spacing=0.95)
glyphmap = {"1": ("✓", TEAL), "0": ("—", "B9C2C9"), "~": ("~", GOLD)}
for i, row in enumerate(rows):
    y = ty0 + rh + 0.06 + i * rh
    band = PALE2 if i % 2 == 0 else WHITE
    box(s, tx0, y, sum(cw), rh, fill=band, line=None, radius=0,
        shape=MSO_SHAPE.RECTANGLE)
    txt(s, tx0 + 0.14, y, cw[0] - 0.2, rh, row[0], size=9.5, bold=True,
        color=INK, anchor=MSO_ANCHOR.MIDDLE)
    for j in range(1, 6):
        g, gc = glyphmap[row[j]]
        txt(s, xacc[j], y, cw[j], rh, g, size=12 if g == "✓" else 11, bold=True,
            color=gc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# AURA column outline
box(s, xacc[5], ty0, cw[5], rh + 0.06 + len(rows) * rh, fill=None, line=TEAL,
    lw=1.75, radius=0.04)
txt(s, 0.5, 5.92, 9.0, 0.26,
    [[("The unit of novelty is the ", {"color": MUTED, "italic": True}),
      ("stack", {"color": TEAL, "bold": True, "italic": True}),
      (": epistemic reasoning + quantum representation + planning, in one clinical loop.",
       {"color": MUTED, "italic": True})]], size=10.5, align=PP_ALIGN.CENTER)
cite(s, "Rajpurkar 2017 · Tu et al., arXiv:2307.14334 (2023) · Cardoso et al., arXiv:2211.02701 (2022) · Luo et al., Briefings in Bioinformatics 23(6) (2022)", y=6.22)
notes(s, "POSITIONING  ·  0:45\n"
         "Fair comparison against the strongest baselines - and credit where "
         "due: CheXNet solved probability, MedPaLM-M can suggest follow-ups in "
         "free text, MONAI has calibration utilities. But no published system "
         "combines decomposed uncertainty, guaranteed coverage, counterfactuals, "
         "information-theoretic test planning, and a quantum representation "
         "layer in one auditable clinical loop. The novelty is not any single "
         "row - it is the column.\n\n"
         "JUDGES SHOULD FEEL: the team knows the literature cold and positions "
         "honestly (the '~' marks buy enormous credibility).\n"
         "ANIMATION: baseline columns appear together; AURA column sweeps in "
         "last with the outline.")

# ==============================================================================
# SLIDE 14 — future roadmap
# ==============================================================================
s = S[13]
kicker(s, "12", "ROADMAP")
title(s, "From one scanner to a quantum clinical network")
purpose(s, "Purpose - a staged, hardware-realistic research program - not science fiction.")

line(s, 0.7, 4.32, 9.3, 4.32, color=TEAL2, wpt=1.75, dash="sysDash")
mile = [
    ("2026 H2", "Clinical pilot", "retrospective multi-site validation on public CXR cohorts", True),
    ("2027", "Quantum federated learning", "hospital-local encoders; only kernel statistics leave the site", False),
    ("2027-28", "Digital-twin priors", "patient-specific progression models seed the belief graph", True),
    ("2028", "Hardware kernels", "similarity search on early fault-tolerant QPUs", False),
    ("2029", "ICU real-time agents", "streaming vitals + imaging, continuous belief updates", True),
]
mx = 0.9
for datev, name, desc, up in mile:
    glyph_circle(s, mx + 0.55, 4.32, 0.16, "", fill=NEON, ring=TEAL2)
    yb = 2.98 if up else 4.66
    box(s, mx, yb, 1.72, 1.06, fill=PALE2, line=BORDER)
    txt(s, mx + 0.12, yb + 0.08, 1.5, 0.2, datev, size=8.5, bold=True, color=TEAL, spc=120)
    txt(s, mx + 0.12, yb + 0.28, 1.5, 0.34, name, size=9.5, bold=True, color=INK,
        line_spacing=0.95)
    txt(s, mx + 0.12, yb + 0.62, 1.5, 0.42, desc, size=7.2, color=MUTED,
        line_spacing=0.95)
    line(s, mx + 0.55, 4.32 + (-0.08 if up else 0.08), mx + 0.55,
         (yb + 1.06 if up else yb), color=BORDER, wpt=1.0)
    mx += 1.74
txt(s, 0.5, 6.0, 9.0, 0.26,
    "Every stage runs on simulators today and inherits hardware as it matures - no stage depends on a miracle.",
    size=10, color=TEAL, italic=True, align=PP_ALIGN.CENTER) if False else None
cite(s, "QFL: Chehimi & Saad, ICASSP (2022)  ·  digital twins: Laubenbacher et al., npj Digital Medicine 5, 64 (2022)")
notes(s, "ROADMAP  ·  0:35\n"
         "The roadmap is staged to be hardware-realistic. Next: retrospective "
         "validation on public chest X-ray cohorts. 2027: quantum federated "
         "learning - encoders stay inside each hospital, only kernel statistics "
         "cross the wall, which is a natural privacy fit. Then digital-twin "
         "priors feeding the belief graph, hardware kernels when early "
         "fault-tolerant machines arrive, and streaming ICU agents. Note what "
         "is absent: no stage requires a quantum miracle - everything runs on "
         "simulators today and inherits hardware as it matures.\n\n"
         "JUDGES SHOULD FEEL: a fundable research program with judgment, not a "
         "wish list.\n"
         "ANIMATION: timeline draws left to right; milestone cards rise in "
         "sequence.")

# ==============================================================================
# SLIDE 15 — impact
# ==============================================================================
s = S[14]
kicker(s, "13", "IMPACT")
title(s, "Who feels this, and how")
purpose(s, "Purpose - translate architecture into outcomes for every stakeholder.")

imp = [
    ("HOSPITALS", "Σ", "worklists triaged by uncertainty - critical, confident cases fast-tracked; reading backlog shrinks"),
    ("DOCTORS", "ψ", "a second reader that shows its work and defers when unsure - trust by design, not decree"),
    ("PATIENTS", "❤", "EIG-planned workups mean fewer unnecessary scans, less radiation, faster answers"),
    ("RESEARCHERS", "∫", "an open uncertainty benchmark + full provenance trails - every posterior reproducible"),
    ("THE SYSTEM", "Δ", "diagnostic error touches ~12 M US adults yearly - calibrated triage attacks the costliest failure in medicine"),
]
positions = [(0.5, 2.52), (3.55, 2.52), (6.6, 2.52), (2.0, 4.42), (5.1, 4.42)]
for (name, g, desc), (x, y) in zip(imp, positions):
    wdt = 2.9
    box(s, x, y, wdt, 1.68, fill=PALE2, line=BORDER)
    glyph_circle(s, x + 0.42, y + 0.40, 0.44, g, fill=PALE, ring=TEAL2,
                 gcolor=TEAL, gsize=14)
    txt(s, x + 0.74, y + 0.26, 2.0, 0.28, name, size=10.5, bold=True, color=INK,
        spc=140)
    txt(s, x + 0.22, y + 0.74, wdt - 0.44, 0.86, desc, size=9, color=MUTED,
        line_spacing=1.08)
cite(s, "Singh et al., BMJ Quality & Safety 23, 727-731 (2014)  ·  WHO Global Atlas of Medical Devices (2022)")
notes(s, "IMPACT  ·  0:35\n"
         "Concretely. Hospitals get uncertainty-triaged worklists - the "
         "confident criticals jump the queue. Doctors get a second reader that "
         "shows its work and knows when to stay silent. Patients get workups "
         "planned by information gain - fewer scans, less radiation, faster "
         "answers. Researchers get reproducible posteriors and an open "
         "benchmark. And the system: diagnostic error touches twelve million "
         "American adults a year - India's numbers are larger and less "
         "measured. Calibrated triage is aimed at the single costliest failure "
         "mode in medicine.\n\n"
         "JUDGES SHOULD FEEL: heart after math - the equations were for "
         "something.\n"
         "ANIMATION: five cards cascade; keep it quick, 0.2 s each.")

# ==============================================================================
# SLIDE 16 — closing (dark)
# ==============================================================================
s = S[15]
dark_bg(s)
rings(s, 5.0, 3.35, [1.35, 2.15, 2.95], alphas=(22, 12, 6))
txt(s, 1.0, 2.35, 8.0, 2.0,
    [[("Artificial intelligence predicts.", {"color": "9FB3B1", "size": 20})],
     [("Clinical intelligence reasons.", {"color": WHITE, "size": 24})],
     [("Quantum intelligence understands.", {"color": NEONHI, "size": 28, "bold": True})]],
    align=PP_ALIGN.CENTER, line_spacing=1.35)
line(s, 4.2, 5.28, 5.8, 5.28, color=NEON, wpt=1.0, alpha=45)
txt(s, 1.0, 5.44, 8.0, 0.62,
    [[("AURA  ·  live demo on request", {"color": WHITE, "size": 11, "bold": True, "spc": 180})],
     [("Sivaguru R.M  ·  Rajalakshmi Institute of Technology  ·  sivagurumurugan1@gmail.com",
       {"color": DIM, "size": 9.5, "spc": 120})]],
    align=PP_ALIGN.CENTER, line_spacing=1.5)
notes(s, "CLOSING  ·  0:25\n"
         "Three sentences. Artificial intelligence predicts - we have had that "
         "for a decade. Clinical intelligence reasons - that is AURA's "
         "architecture. Quantum intelligence understands - a richer geometry "
         "for medical evidence, claimed honestly and benchmarked ruthlessly. "
         "The system is live - I would love to show you. Thank you.\n\n"
         "JUDGES SHOULD FEEL: goosebumps, then the urge to ask for the demo.\n"
         "ANIMATION: three lines fade in on successive clicks, 0.6 s each - the "
         "only slide with deliberate dramatic pacing.")

# ==============================================================================
# SLIDE 17 — appendix A: mathematics
# ==============================================================================
s = S[16]
kicker(s, "A1", "TECHNICAL APPENDIX - MATHEMATICAL CORE")
title(s, "The four equations under the hood", size=24)

eqcards = [
    ("DENSITY-MATRIX EMBEDDING & QUANTUM KERNEL",
     ["ρ(x) = |φ(x)⟩⟨φ(x)|",
      "k(x, x′) = Tr[ ρ(x) ρ(x′) ] = |⟨φ(x)|φ(x′)⟩|²"],
     "kernel = state fidelity; PSD by construction, estimated from shot statistics"),
    ("VARIATIONAL QUANTUM CLASSIFIER (FUSION)",
     ["fθ(x) = ⟨0| Uφ†(x) W†(θ) M W(θ) Uφ(x) |0⟩",
      "∂f/∂θᵢ = ½ [ f(θᵢ+π/2) − f(θᵢ−π/2) ]"],
     "trained by parameter-shift - exact gradients, no backprop through the circuit"),
    ("BAYESIAN MEASUREMENT UPDATE",
     ["ρ′ = Eₒ ρ Eₒ† / Tr[ Eₒ ρ Eₒ† ]"],
     "quantum state update generalizes Bayes' rule (Leifer & Spekkens 2013) - evidence as measurement"),
    ("CONFORMAL + PLANNING LAYER",
     ["C(x) = { y : s(x,y) ≤ q̂₁₋α },   P(y ∈ C) ≥ 1−α",
      "a* = argmaxₐ { H[p(d)] − EₒH[p(d|o)] − λ c(a) }"],
     "finite-sample coverage + Lindley-optimal next test, on top of any posterior"),
]
positions = [(0.5, 2.30), (5.05, 2.30), (0.5, 4.24), (5.05, 4.24)]
for (name, eqs, sub), (x, y) in zip(eqcards, positions):
    box(s, x, y, 4.45, 1.80, fill=PALE2, line=BORDER)
    txt(s, x + 0.2, y + 0.13, 4.1, 0.22, name, size=8, color=TEAL, bold=True, spc=140)
    yy = y + 0.40
    for eq in eqs:
        txt(s, x + 0.2, yy, 4.08, 0.3, eq, size=11.5, color=INK, bold=True,
            font="Cambria")
        yy += 0.36
    txt(s, x + 0.2, y + 1.42, 4.08, 0.34, sub, size=8, color=MUTED, italic=True,
        line_spacing=0.95)
cite(s, "Leifer & Spekkens, Phys. Rev. A 88, 052130 (2013)  ·  Mitarai et al., Phys. Rev. A 98, 032309 (2018) - parameter shift  ·  Schuld, arXiv:2101.11020 (2021)", y=6.18)
notes(s, "APPENDIX A - keep in reserve for the quantum jury.\n"
         "Walk any equation on request: kernel PSD-ness from fidelity; "
         "parameter-shift for exact gradients on hardware; the Leifer-Spekkens "
         "view of Bayesian updating as CPTP maps; conformal quantile on the "
         "calibration split. If asked 'why not amplitude encoding?': amplitude "
         "encoding needs normalized inputs and deep state-prep circuits; angle "
         "encoding on 8 channels is NISQ-shallow and keeps features "
         "interpretable per qubit.")

# ==============================================================================
# SLIDE 18 — appendix B: implementation + benchmark
# ==============================================================================
s = S[17]
kicker(s, "A2", "TECHNICAL APPENDIX - WHAT RUNS TODAY")
title(s, "Built, not proposed", size=24)

box(s, 0.5, 2.35, 4.4, 3.55, fill=PALE2, line=BORDER)
txt(s, 0.72, 2.50, 4.0, 0.24, "IMPLEMENTED STACK  (offline-capable)", size=9,
    color=TEAL, bold=True, spc=180)
impl = [
    ("Quantum fusion", "PennyLane statevector VQC - 8-qubit angle encoding, entangling layers"),
    ("Classical twin", "Bayesian product-of-experts fusion behind the same interface"),
    ("Safety engine", "deep ensembles + MC-dropout, temperature scaling, conformal sets, energy-score OOD, abstention"),
    ("Explainability", "occlusion saliency, Shapley-style evidence attribution, counterfactual sensitivity"),
    ("Recommender", "expected-information-gain next-test ranking per cost/risk"),
    ("Platform", "FastAPI gateway, async event bus, SQLite audit ledger, doctor dashboard SPA"),
]
iy = 2.82
for name, desc in impl:
    txt(s, 0.72, iy, 1.35, 0.5, name, size=8.8, bold=True, color=INK)
    txt(s, 2.12, iy, 2.62, 0.5, desc, size=8, color=MUTED, line_spacing=0.95)
    iy += 0.52

box(s, 5.15, 2.35, 4.35, 2.0, fill=PALE, line=TEAL2, lw=1.1)
txt(s, 5.36, 2.50, 4.0, 0.24, "THE BENCHMARK LAW", size=9, color=TEAL, bold=True,
    spc=200)
txt(s, 5.36, 2.78, 3.95, 1.0,
    [[("aura_cli bench", {"font": "Courier New", "bold": True, "color": INK, "size": 10.5}),
      ("  runs quantum vs classical twin on identical held-out data:",
       {"color": MUTED, "size": 9})],
     [("accuracy · NLL · ECE · Brier · conformal coverage · set size",
       {"color": INK, "size": 9.5, "bold": True})]],
    line_spacing=1.15)
txt(s, 5.36, 3.80, 3.95, 0.44,
    "If the classical twin wins, the classical twin ships - and we say so on stage.",
    size=9.5, color=TEAL, italic=True, bold=True, line_spacing=1.05)

box(s, 5.15, 4.55, 4.35, 1.35, fill=WHITE, line=BORDER)
txt(s, 5.36, 4.68, 4.0, 0.24, "LIVE DEMO PATH", size=9, color=TEAL, bold=True, spc=200)
demo = [("py -m aura_cli train", "fit fusion + calibration (seconds)"),
        ("py -m aura_cli bench", "quantum vs classical, head-to-head"),
        ("py -m aura_cli serve", "doctor dashboard on :8000")]
dyy = 4.94
for cmd, what in demo:
    txt(s, 5.36, dyy, 2.2, 0.24, cmd, size=9, color=INK, font="Courier New", bold=True)
    txt(s, 7.42, dyy, 2.0, 0.24, what, size=8, color=MUTED, italic=True)
    dyy += 0.30
cite(s, "Full provenance per result: {backend, device, feature-map version, circuit version, seed, shots} pinned to every belief snapshot.")
notes(s, "APPENDIX B - use when judges ask 'is this real?'\n"
         "Everything on the left runs offline on a laptop. The benchmark law is "
         "the differentiator: the quantum path must beat its classical twin on "
         "the same held-out data across six metrics, or the classical twin "
         "serves. Offer the live demo here - train, bench, serve - three "
         "commands, under two minutes.")

# ==============================================================================
# SLIDE 19 — references
# ==============================================================================
s = S[18]
kicker(s, "A3", "REFERENCES")
title(s, "References", size=24)
refs_l = [
    "Havlíček, Córcoles, Temme et al. Supervised learning with quantum-enhanced feature spaces. Nature 567, 209-212 (2019).",
    "Liu, Arunachalam & Temme. A rigorous and robust quantum speed-up in supervised machine learning. Nature Physics 17, 1013-1017 (2021).",
    "Schuld & Killoran. Quantum machine learning in feature Hilbert spaces. Physical Review Letters 122, 040504 (2019).",
    "Mitarai, Negoro, Kitagawa & Fujii. Quantum circuit learning. Physical Review A 98, 032309 (2018).",
    "Cerezo et al. Variational quantum algorithms. Nature Reviews Physics 3, 625-644 (2021).",
    "Leifer & Spekkens. Towards a formulation of quantum theory as a causally neutral theory of Bayesian inference. Physical Review A 88, 052130 (2013).",
    "Kendall & Gal. What uncertainties do we need in Bayesian deep learning for computer vision? NeurIPS (2017).",
]
refs_r = [
    "Guo, Pleiss, Sun & Weinberger. On calibration of modern neural networks. ICML (2017).",
    "Angelopoulos & Bates. Conformal prediction: a gentle introduction. Foundations and Trends in Machine Learning 16(4) (2023).",
    "Lindley. On a measure of the information provided by an experiment. Annals of Mathematical Statistics 27, 986-1005 (1956).",
    "Richens, Lee & Johri. Improving the accuracy of medical diagnosis with causal machine learning. Nature Communications 11, 3923 (2020).",
    "Rajpurkar et al. CheXNet: radiologist-level pneumonia detection on chest X-rays. arXiv:1711.05225 (2017).",
    "Zech et al. Variable generalization performance of a deep learning model to detect pneumonia. PLOS Medicine 15(11) (2018).",
    "Rudin. Stop explaining black box machine learning models for high stakes decisions. Nature Machine Intelligence 1, 206-215 (2019).",
]
for i, r in enumerate(refs_l):
    txt(s, 0.5, 2.35 + i * 0.52, 4.4, 0.5, r, size=8, color=MUTED, line_spacing=0.95)
for i, r in enumerate(refs_r):
    txt(s, 5.1, 2.35 + i * 0.52, 4.4, 0.5, r, size=8, color=MUTED, line_spacing=0.95)
notes(s, "REFERENCES - leave up during Q&A. Every claim on every slide traces "
         "to one of these.")

prs.save("AURA_QuantAThan_2026.pptx")
print("saved AURA_QuantAThan_2026.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
