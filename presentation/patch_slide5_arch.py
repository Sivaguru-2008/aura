# -*- coding: utf-8 -*-
"""Slide 5 — dark 'cyber-medical' HYBRID MACHINE visualizer (maxed-out).

3 layered tiers, cyan classical services, a glowing magenta QPU module docked into
the FUSION bay (multi-halo glow + online status + ghost classical twin), a nested
fidelity-kernel readout, HUD corner brackets, a wired power spine, circuit traces,
and a tech-stack footer bar. Light chrome (kicker/title/logos/footer) preserved.

Idempotent (clears + rebuilds).  Run:  py -X utf8 patch_slide5_arch.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"

INK, MUTED = "0F172A", "5B6672"
TEAL, NEON, NEONHI = "0E7490", "2DD4BF", "5EEAD4"
MAG, MAGHI, MAGDK = "D946EF", "F4B6FB", "3E1642"
LIME = "5EF08A"
WHITE, DARKBG, PANEL, DIM = "FFFFFF", "05080D", "0C1720", "8FB6B3"

def _rgb(h): return RGBColor.from_string(h)
def _strip(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None: sp._element.remove(st)
def _alpha(fill_elm, pct):
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is not None:
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))

def box(s, x, y, w, h, fill=PANEL, line=None, lw=0.75, radius=0.10,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None, fill_alpha=None, line_alpha=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip(sp); sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None: sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha is not None: _alpha(sp.fill.fore_color._xFill, fill_alpha)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
        ln = sp.line._get_or_add_ln()
        if dash is not None:
            ln.insert(list(ln).index(ln.find(qn("a:solidFill"))) + 1,
                      ln.makeelement(qn("a:prstDash"), {"val": dash}))
        if line_alpha is not None:
            sf = ln.find(qn("a:solidFill"))
            if sf is not None: _alpha(sf, line_alpha)
    return sp

def dot(s, cx, cy, d, color, alpha=None):
    box(s, cx - d/2, cy - d/2, d, d, fill=color, line=None, shape=MSO_SHAPE.OVAL, fill_alpha=alpha)

def txt(s, x, y, w, h, content, size=8, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial", spc=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = [[(t, {})] for t in content.split("\n")] if isinstance(content, str) else content
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = opts.get("font", font)
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None: r._r.get_or_add_rPr().set("spc", str(sv))
    return tb

def conn(s, x1, y1, x2, y2, color=NEON, wpt=1.0, alpha=None, arrow=False, dash=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _strip(c); c.shadow.inherit = False
    c.line.color.rgb = _rgb(color); c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    if dash: ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    if arrow: ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if alpha is not None:
        sf = ln.find(qn("a:solidFill"))
        if sf is not None: _alpha(sf, alpha)
    return c

def bracket(s, x, y, sx, sy, color=NEON):
    box(s, x, y, sx*0.34, 0.03, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE, fill_alpha=70)
    box(s, x if sx > 0 else x, y, 0.03, sy*0.34, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE, fill_alpha=70)

def slab(s, x, y, w, h, accent=NEON):
    box(s, x+0.07, y+0.07, w, h, fill=accent, line=None, radius=0.08, fill_alpha=9)   # depth
    box(s, x, y, w, h, fill=PANEL, line=accent, lw=1.0, radius=0.08, line_alpha=60)
    box(s, x, y, 0.05, h, fill=accent, line=None, radius=0.0, fill_alpha=55)          # accent edge

def svc(s, x, y, w, h, label):
    box(s, x, y, w, h, fill=PANEL, line=NEON, lw=1.0, radius=0.14, line_alpha=75)
    box(s, x+0.10, y+0.05, w-0.20, 0.035, fill=NEONHI, line=None, shape=MSO_SHAPE.RECTANGLE, fill_alpha=70)
    txt(s, x, y+0.03, w, h-0.03, label, size=7.3, color=NEONHI, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dot(s, x+w-0.11, y+h-0.11, 0.055, LIME)

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__); sys.exit(1)

s = prs.slides[4]
for sh in list(s.shapes):
    top = sh.top/914400.0
    if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER: continue
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and top < 1.5: continue
    sh._element.getparent().remove(sh._element)

# light chrome
txt(s, 0.5, 1.24, 9.0, 0.26,
    [[("04", {"color": MUTED}), ("  ·  ", {"color": MUTED}), ("ARCHITECTURE & STACK", {})]],
    size=10.5, color=TEAL, bold=True, spc=280)
txt(s, 0.5, 1.50, 9.2, 0.52, "Modular 3-tier — quantum as a swappable service",
    size=24, color=INK, bold=True, spc=-20)

# ===== dark panel + HUD frame =====
box(s, 0.4, 2.06, 9.2, 3.82, fill=DARKBG, line=NEON, lw=1.1, radius=0.04, line_alpha=30)
box(s, 0.4, 2.06, 9.2, 3.82, fill=None, line=NEON, lw=0.5, radius=0.04, line_alpha=12)
for cx, cy, sx, sy in [(0.5,2.16,1,1),(9.5,2.16,-1,1),(0.5,5.78,1,-1),(9.5,5.78,-1,-1)]:
    box(s, cx-(0.34 if sx<0 else 0), cy-(0.34 if sy<0 else 0)+ (0 if sy>0 else 0.31), 0.34, 0.028,
        fill=NEONHI, line=None, shape=MSO_SHAPE.RECTANGLE, fill_alpha=75)
    box(s, cx-(0.028 if sx<0 else 0), cy-(0.34 if sy<0 else 0), 0.028, 0.34,
        fill=NEONHI, line=None, shape=MSO_SHAPE.RECTANGLE, fill_alpha=75)

# faint circuit traces (texture, top-right empty band)
for yy in (3.60, 3.78):
    conn(s, 6.9, yy, 8.9, yy, color=NEON, wpt=0.8, alpha=16)
    dot(s, 8.9, yy, 0.06, NEON, alpha=40)
conn(s, 7.4, 3.60, 7.4, 3.78, color=NEON, wpt=0.8, alpha=16)

txt(s, 0.72, 2.20, 6.0, 0.24,
    [[("◇ THE HYBRID MACHINE   ", {"color": NEONHI, "bold": True, "size": 10.5, "spc": 90}),
      ("— quantum docks into the fusion bay", {"color": DIM, "size": 8, "italic": True})]])
txt(s, 6.5, 2.22, 2.9, 0.22,
    [[("● ", {"color": NEON, "size": 9}), ("classical   ", {"color": DIM, "size": 7.5}),
      ("● ", {"color": MAG, "size": 9}), ("quantum", {"color": DIM, "size": 7.5})]],
    align=PP_ALIGN.RIGHT)

# ===== power spine (left, wiring the tiers) =====
conn(s, 0.60, 2.82, 0.60, 4.97, color=NEON, wpt=1.4, alpha=50)
for ny in (2.80, 3.90, 4.97):
    dot(s, 0.60, ny, 0.13, NEON, alpha=22); dot(s, 0.60, ny, 0.075, NEONHI)

# ===== PRESENTATION tier =====
slab(s, 0.8, 2.54, 8.4, 0.52)
txt(s, 1.02, 2.54, 2.0, 0.52, "PRESENTATION", size=8.3, color=NEON, bold=True, anchor=MSO_ANCHOR.MIDDLE, spc=60)
txt(s, 2.95, 2.54, 6.1, 0.52, "doctor dashboard · worklist · case console · report signing",
    size=8, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ===== APPLICATION tier =====
slab(s, 0.8, 3.18, 8.4, 1.46)
txt(s, 1.02, 3.24, 2.0, 0.22, "APPLICATION", size=8.3, color=NEON, bold=True, spc=60)
txt(s, 2.6, 3.24, 6.4, 0.22, "FastAPI gateway · async event bus · typed Pydantic contracts",
    size=7.4, color=DIM)

tx, tw, tg = 1.0, 0.92, 0.09
xs = [tx + i*(tw+tg) for i in range(8)]
labels = ["vision", "fusion", "safety", "explain", "recommend", "report", "memory", "OPEN"]
for i, nm in enumerate(labels):
    x = xs[i]
    if nm == "OPEN":
        box(s, x, 4.06, tw, 0.48, fill=DARKBG, line=DIM, lw=1.0, radius=0.14, dash="sysDash", line_alpha=55)
        txt(s, x, 4.06, tw, 0.48, [[("OPEN BAY", {"color": DIM, "bold": True, "size": 6.8})],
                                   [("add svc", {"color": MUTED, "size": 5.8})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    elif nm == "fusion":
        pass  # drawn below with the QPU
    else:
        svc(s, x, 4.06, tw, 0.48, nm)

# ---- FUSION bay + ghost classical twin + QPU module ----
fx = xs[1]
box(s, fx+0.12, 4.12, tw, 0.44, fill=None, line=NEON, lw=1.0, radius=0.14, dash="sysDash", line_alpha=32)  # twin ghost
txt(s, fx+0.12, 4.40, tw, 0.16, "classical twin", size=5.4, color=NEON, italic=True, align=PP_ALIGN.CENTER)
box(s, fx, 4.06, tw, 0.48, fill="14232C", line=MAG, lw=1.4, radius=0.14)                                   # dock bay
txt(s, fx, 4.06, tw, 0.48, "fusion", size=7.4, color=MAGHI, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
dot(s, fx+tw-0.11, 4.17, 0.055, MAGHI)

# multi-halo glow + pulse ring
box(s, fx-0.30, 3.30, tw+0.60, 0.86, fill=MAG, line=None, radius=0.30, fill_alpha=5)
box(s, fx-0.22, 3.36, tw+0.44, 0.72, fill=MAG, line=None, radius=0.30, fill_alpha=8)
box(s, fx-0.13, 3.42, tw+0.26, 0.60, fill=MAG, line=None, radius=0.28, fill_alpha=13)
box(s, fx-0.30, 3.28, tw+0.60, 0.90, fill=None, line=MAGHI, lw=1.0, radius=0.30, shape=MSO_SHAPE.OVAL, line_alpha=30)
box(s, fx-0.02, 3.50, tw+0.04, 0.46, fill=MAGDK, line=MAGHI, lw=1.4, radius=0.16)                          # core
txt(s, fx-0.02, 3.51, tw+0.04, 0.30,
    [[("◆ QPU", {"color": MAGHI, "bold": True, "size": 8})],
     [("8-qubit VQC", {"color": "EBC4F2", "size": 6})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
dot(s, fx+0.16, 3.87, 0.05, LIME); txt(s, fx+0.22, 3.79, 0.7, 0.14, "ONLINE", size=5, color=LIME, bold=True)
conn(s, fx+tw/2, 3.96, fx+tw/2, 4.06, color=MAGHI, wpt=1.8, arrow=True)                                    # dock plug

# nested fidelity-kernel readout
box(s, fx+tw+0.20, 3.50, 3.55, 0.50, fill=MAG, line=MAGHI, lw=0.9, radius=0.12, fill_alpha=7, line_alpha=45)
txt(s, fx+tw+0.34, 3.52, 3.3, 0.46,
    [[("k(x,x′) = |⟨φ(x)|φ(x′)⟩|²", {"color": MAGHI, "size": 10, "font": "Consolas", "bold": True})],
     [("fidelity kernel · classical twin sits behind the same interface",
       {"color": DIM, "size": 6.4, "italic": True})]], line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)

# ===== DATA tier =====
slab(s, 0.8, 4.72, 8.4, 0.50)
txt(s, 1.02, 4.72, 1.6, 0.50, "DATA", size=8.3, color=NEON, bold=True, anchor=MSO_ANCHOR.MIDDLE, spc=60)
txt(s, 2.45, 4.72, 6.7, 0.50, "audit ledger · case memory · model registry · provenance {seed · circuit · shots}",
    size=7.6, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ===== TECH STACK footer bar =====
box(s, 0.8, 5.30, 8.4, 0.34, fill="0A141B", line=NEON, lw=0.75, radius=0.10, line_alpha=45)
stack = ["Python 3", "FastAPI", "PennyLane", "NumPy", "scikit-learn", "SciPy", "SQLite", "Pydantic"]
runs = [("▚ STACK   ", {"color": NEON, "bold": True, "size": 7.4, "spc": 120})]
for i, it in enumerate(stack):
    runs.append((it, {"color": WHITE, "size": 7.4}))
    if i < len(stack)-1: runs.append(("   ·   ", {"color": NEON, "size": 7.4}))
txt(s, 1.0, 5.30, 8.1, 0.34, [runs], anchor=MSO_ANCHOR.MIDDLE)

# caption (light, outside panel)
txt(s, 0.5, 5.96, 9.0, 0.22,
    "Offline-capable · CPU-only · no PHI leaves the site · bus → Redis Streams, DB → Postgres in production",
    size=7.6, color=MUTED, italic=True)

s.notes_slide.notes_text_frame.text = (
    "ARCHITECTURE & STACK  ·  0:45\n"
    "AURA is a machine with three layered tiers: a data foundation, an application "
    "tier of independent services behind one FastAPI gateway, and the doctor-facing "
    "presentation tier. Modularity is the point: every service is a slot. Our "
    "quantum VQC — the glowing magenta QPU — docks into the FUSION bay, and a "
    "classical twin sits behind the exact same interface, so whichever wins the "
    "benchmark ships with zero code change. Open bays take new services the same "
    "way. Cyan is classical infrastructure; magenta is the quantum accelerator. "
    "It all runs offline, CPU-only, no PHI leaving the site.\n\n"
    "JUDGES SHOULD FEEL: disciplined engineering — quantum is a swappable part, not a bet.")

prs.save(DECK)
print("saved", DECK, "— Slide 5 maxed-out hybrid-machine visualizer")
