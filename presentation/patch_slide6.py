# -*- coding: utf-8 -*-
"""In-place patch of Slide 6: replace the '24-hour plan' left column with a
QUANTUM STACK explainer (what's used / used efficiently / why it helps).

Non-destructive: keeps the benchmark-receipts terminal, the hand-added logos,
the header/footer bars, and the honesty-contract citation. Edits only the
kicker+title text and swaps the left content column.

Run:  py -X utf8 patch_slide6.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"

# ---- palette (from build_pitch.py) -----------------------------------------
INK, MUTED = "0F172A", "5B6672"
TEAL, TEAL2, NEON, NEONHI = "0E7490", "0D9488", "2DD4BF", "5EEAD4"
PALE, PALE2, BORDER = "ECF8F7", "F6FAFA", "C6E4E1"
WHITE, DARKBG, DIM = "FFFFFF", "05090E", "7FA8A5"
WARN, GOLD, GRAY = "C05B52", "B98A2F", "AEB8C2"

def _rgb(h): return RGBColor.from_string(h)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)

def _set_alpha(fill_elm, pct):
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is not None:
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))

def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.10,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None, fill_alpha=None,
        line_alpha=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp); sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha is not None:
            _set_alpha(sp.fill.fore_color._xFill, fill_alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
        ln = sp.line._get_or_add_ln()
        if line_alpha is not None:
            sf = ln.find(qn("a:solidFill"))
            if sf is not None: _set_alpha(sf, line_alpha)
    return sp

def txt(s, x, y, w, h, content, size=14, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial", spc=None,
        wrap=True, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(content, str):
        paras = [[(t, {})] for t in content.split("\n")]
    else:
        paras = [([item] if isinstance(item, str) else item) if not isinstance(item, str)
                 else [(item, {})] for item in content]
        paras = []
        for item in content:
            paras.append([(item, {})] if isinstance(item, str) else item)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.name = opts.get("font", font)
            f.color.rgb = _rgb(opts.get("color", color))
            sp_v = opts.get("spc", spc)
            if sp_v is not None:
                r._r.get_or_add_rPr().set("spc", str(sp_v))
    return tb

def chip(s, x, y, w, h, label, fill=PALE, lncol=TEAL2, tcol=TEAL, size=8.5, bold=True):
    box(s, x, y, w, h, fill=fill, line=lncol, lw=0.75, radius=0.5)
    txt(s, x, y, w, h, label, size=size, color=tcol, bold=bold,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def in_(emu): return emu / 914400.0

# ---- load ------------------------------------------------------------------
try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__)
    sys.exit(1)

s = prs.slides[5]  # slide 6

# 1) retitle kicker + title in place -----------------------------------------
for sh in s.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if "IMPLEMENTATION & PROOF" in t:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "IMPLEMENTATION" in r.text:
                    r.text = "QUANTUM STACK"
    elif "24-hour plan" in t:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = "What's quantum — and why it earns its seat"

# 2) delete the left build-plan column (keep pictures / terminal / bars) ------
removed = []
for sh in list(s.shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        continue
    L, T = in_(sh.left), in_(sh.top)
    if L < 4.90 and 2.05 < T < 5.95:
        removed.append((round(L, 2), round(T, 2),
                        (sh.text_frame.text[:34] if sh.has_text_frame else "<%s>" % sh.shape_type)))
        sh._element.getparent().remove(sh._element)
print("removed %d left-column shapes:" % len(removed))
for r in removed:
    print("   ", r)

# 3) draw the new left column ------------------------------------------------
CODE = TEAL2  # colour for code / formula fragments

# Panel A -- WHAT'S USED (the circuit)
box(s, 0.5, 2.14, 4.35, 1.90, fill=PALE2, line=BORDER)
txt(s, 0.68, 2.24, 4.0, 0.22, "WHAT'S USED  ·  THE CIRCUIT (SHIPS TODAY)",
    size=8, color=TEAL, bold=True, spc=140)
steps = [
    ("ENCODE", [("RY(π·xᵢ)", {"color": CODE, "bold": True, "font": "Consolas", "size": 8.5}),
                ("  ·  1 qubit per evidence channel (8)", {"size": 8.2, "color": INK})]),
    ("ENTANGLE", [("3×(RY·RZ + CNOT ring)", {"color": CODE, "bold": True, "font": "Consolas", "size": 8.5}),
                  ("  →  higher-order interactions, 2⁸ space", {"size": 8.2, "color": INK})]),
    ("READ", [("⟨Zᵢ⟩", {"color": CODE, "bold": True, "font": "Consolas", "size": 8.5}),
              ("  →  linear head  →  6-dx calibrated posterior", {"size": 8.2, "color": INK})]),
]
yy = 2.56
for lab, desc in steps:
    txt(s, 0.68, yy, 1.02, 0.24, lab, size=8, bold=True, color=TEAL, spc=60)
    txt(s, 1.62, yy, 3.15, 0.30, [desc], wrap=True, line_spacing=1.0)
    yy += 0.42
for cx, w, lab in [(0.68, 0.82, "8 qubits"), (1.56, 0.92, "48 params"),
                   (2.54, 0.86, "512 shots"), (3.46, 1.32, "CPU · default.qubit")]:
    chip(s, cx, 3.62, w, 0.30, lab, fill=PALE, lncol=TEAL2, tcol=TEAL, size=7.4)

# Panel B -- USED EFFICIENTLY / WHY IT HELPS
box(s, 0.5, 4.14, 4.35, 1.78, fill=PALE2, line=BORDER)
txt(s, 0.68, 4.24, 4.0, 0.22, "USED EFFICIENTLY  ·  WHY IT HELPS",
    size=8, color=TEAL, bold=True, spc=140)
bullets = [
    [("Compact.  ", {"bold": True, "color": TEAL, "size": 8.2}),
     ("O(n·layers) params where a classical interaction model needs exponentially many terms.",
      {"size": 8.2, "color": INK})],
    [("Uncertainty for free.  ", {"bold": True, "color": TEAL, "size": 8.2}),
     ("Shot variance (1−⟨Z⟩²)/N → posterior_std the safety layer already consumes.",
      {"size": 8.2, "color": INK})],
    [("Never pixels.  ", {"bold": True, "color": TEAL, "size": 8.2}),
     ("8-dim compressed evidence — tractable now, hardware-credible later. Async, off the <300 ms read path.",
      {"size": 8.2, "color": INK})],
]
by = 4.54
for b in bullets:
    box(s, 0.68, by + 0.04, 0.05, 0.32, fill=TEAL2, line=None, radius=0.0,
        shape=MSO_SHAPE.RECTANGLE)
    txt(s, 0.86, by, 3.92, 0.42, [b], wrap=True, line_spacing=1.02, size=8.2)
    by += 0.44

# 4) refresh speaker notes ----------------------------------------------------
s.notes_slide.notes_text_frame.text = (
    "QUANTUM STACK  ·  1:00\n"
    "What is actually quantum: an 8-qubit variational circuit in PennyLane. Each "
    "of the 8 evidence channels is angle-encoded, RY(pi x); three layers of "
    "trainable rotations plus a CNOT ring entangle them, so the model represents "
    "higher-order interactions between evidence sources in a 2^8 space with only "
    "48 parameters. We read Pauli-Z expectations into a linear head for a "
    "calibrated 6-diagnosis posterior. Efficient by design: 48 params, CPU-only, "
    "and finite-shot measurement variance gives uncertainty for free that the "
    "safety layer consumes. We never put pixels in a circuit - only the 8-dim "
    "compressed evidence, which is exactly why this is tractable on a simulator "
    "today and credible on hardware tomorrow. Right panel is the receipt: the "
    "quantum path beats its classical twin on the same held-out data, and if it "
    "did not, the twin would ship. Every claim is tiered in QUANTUM_STACK.md.")

prs.save(DECK)
print("\nsaved", DECK)
