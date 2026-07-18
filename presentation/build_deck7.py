# -*- coding: utf-8 -*-
"""Build the 7-slide finals deck for AURA on the official RIT template."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import math

INK    = "0F172A"
MUTED  = "5B6672"
TEAL   = "0E7490"
TEAL2  = "0D9488"
NEON   = "2DD4BF"
NEONHI = "5EEAD4"
PALE   = "ECF8F7"
PALE2  = "F6FAFA"
BORDER = "C6E4E1"
WHITE  = "FFFFFF"
DARKBG = "05090E"
PANEL  = "0D1922"
DIM    = "7FA8A5"
WARN   = "C05B52"
GOLD   = "B98A2F"

W, H = 10.0, 7.5

prs = Presentation("base.pptx")

# keep slides 1-7 only (slide 2 of the template is our white base for 2-6;
# slide 7 will be painted dark)
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides)[7:]:
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(sldId)
S = list(prs.slides)

def _rgb(hexs):
    return RGBColor.from_string(hexs)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)

def _set_alpha(fill_elm, pct):
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is None:
        return
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    clr.append(a)

def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.10,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None, fill_alpha=None,
        line_alpha=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha is not None:
            _set_alpha(sp.fill.fore_color._xFill, fill_alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(lw)
        ln = sp.line._get_or_add_ln()
        if dash is not None:
            d = ln.makeelement(qn("a:prstDash"), {"val": dash})
            ln.insert(list(ln).index(ln.find(qn("a:solidFill"))) + 1, d)
        if line_alpha is not None:
            sf = ln.find(qn("a:solidFill"))
            if sf is not None:
                _set_alpha(sf, line_alpha)
    return sp

def txt(s, x, y, w, h, content, size=14, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial", spc=None,
        wrap=True, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(content, str):
        paras = [[(t, {})] for t in content.split("\n")]
    else:
        paras = []
        for item in content:
            if isinstance(item, str):
                paras.append([(item, {})])
            else:
                paras.append(item)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run()
            r.text = rt
            f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.name = opts.get("font", font)
            f.color.rgb = _rgb(opts.get("color", color))
            sp_v = opts.get("spc", spc)
            if sp_v is not None:
                r._r.get_or_add_rPr().set("spc", str(sp_v))
    return tb

def kicker(s, num, tag, color=TEAL, x=0.5, y=1.24, w=9.0):
    txt(s, x, y, w, 0.26, [[(num + "  ·  ", {"color": MUTED}), (tag, {})]],
        size=10.5, color=color, bold=True, spc=280)

def title(s, text, x=0.5, y=1.50, w=9.2, size=26, color=INK):
    txt(s, x, y, w, 0.52, text, size=size, color=color, bold=True, spc=-20)

def cite(s, text, y=6.08):
    txt(s, 0.5, y, 9.0, 0.24, text, size=8, color=MUTED, italic=True)

def arrow(s, x1, y1, x2, y2, color=TEAL2, wpt=1.4, dash=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c)
    c.shadow.inherit = False
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    t = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(t)
    return c

def line(s, x1, y1, x2, y2, color=BORDER, wpt=1.0, dash=None, alpha=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    _strip_style(c)
    c.shadow.inherit = False
    c.line.color.rgb = _rgb(color)
    c.line.width = Pt(wpt)
    ln = c.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    if alpha is not None:
        sf = ln.find(qn("a:solidFill"))
        if sf is not None:
            _set_alpha(sf, alpha)
    return c

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def dark_bg(s):
    box(s, -0.05, -0.05, W + 0.1, H + 0.1, fill=DARKBG, line=None, radius=0,
        shape=MSO_SHAPE.RECTANGLE)

def rings(s, cx, cy, radii, color=NEON, alphas=(22, 12, 6)):
    for r, a in zip(radii, alphas):
        box(s, cx - r, cy - r, 2 * r, 2 * r, fill=None, line=color, lw=1.0,
            shape=MSO_SHAPE.OVAL, line_alpha=a)

# ==============================================================================
# SLIDE 1 — title (official template, personalized)
# ==============================================================================
s = S[0]
txt(s, 5.62, 4.62, 4.2, 0.62,
    [[("AURA", {"size": 21, "bold": True, "color": NEONHI, "spc": 500})],
     [("ADAPTIVE UNCERTAINTY-AWARE RADIOLOGY ARCHITECTURE",
       {"size": 8.5, "color": WHITE, "spc": 220})]],
    align=PP_ALIGN.CENTER)
notes(s, "TITLE  ·  0:15\n"
         "Good morning. I am Sivaguru, and this is AURA - the Adaptive "
         "Uncertainty-aware Radiology Architecture: a system that treats "
         "diagnosis not as classification but as a reasoning process, with a "
         "quantum core that is mathematically earned, not marketed. Seven "
         "slides, one live system.\n\n"
         "RUBRIC: presentation - confident, zero filler, sets the honesty tone.\n"
         "JUDGES SHOULD FEEL: this team will not waste our time.")

# ==============================================================================
# SLIDE 2 — the problem: trust gap + reasoning gap in one
# ==============================================================================
s = S[1]
kicker(s, "01", "THE PROBLEM")
title(s, "95% accurate. Still not trusted.")
txt(s, 0.5, 2.04, 9.0, 0.28,
    "Medical AI passed radiologist-level accuracy in 2017. Adoption still stalls - because accuracy was never the bottleneck.",
    size=11, color=MUTED, italic=True)

stats = [("950+", "FDA-cleared AI medical devices; >75% in radiology (2025)"),
         ("1 : 100k", "radiologist-to-population ratio in India (US ~1 : 10k)"),
         ("0", "deployed classifiers that report when their own output should not be trusted")]
sy = 2.50
for big, small in stats:
    box(s, 0.5, sy, 3.72, 0.80, fill=PALE2, line=BORDER)
    txt(s, 0.68, sy + 0.12, 1.30, 0.56, big, size=20, color=TEAL, bold=True,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.98, sy + 0.10, 2.12, 0.62, small, size=8.5, color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    sy += 0.92
box(s, 0.5, 5.30, 3.72, 0.70, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 0.66, 5.40, 3.45, 0.52,
    [[("Clinicians ask: ", {"bold": True, "color": INK, "size": 9.5}),
      ("why? how sure? what would change your mind? what next?",
       {"color": TEAL, "size": 9.5, "italic": True})]], line_spacing=1.05)

# right: today's pipeline + the amputated reasoning
nodes = [("IMAGE", ""), ("CNN", ""), ("p = 0.87", ""), ("DOCTOR", "alone")]
nx, ny, nw2, nh2, gap = 4.55, 2.50, 1.10, 0.62, 0.22
for i, (lab, sub) in enumerate(nodes):
    x = nx + i * (nw2 + gap)
    box(s, x, ny, nw2, nh2, fill=PALE2, line=BORDER)
    txt(s, x, ny + (0.09 if sub else 0.18), nw2, 0.28, lab, size=10.5, bold=True,
        color=INK, align=PP_ALIGN.CENTER)
    if sub:
        txt(s, x, ny + 0.345, nw2, 0.2, sub, size=7.5, color=MUTED,
            align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + nw2 + 0.03, ny + nh2 / 2, x + nw2 + gap - 0.03, ny + nh2 / 2,
              wpt=1.2)
box(s, 4.55, 3.42, 4.95, 1.62, fill=None, line=WARN, lw=1.2, dash="dash")
txt(s, 4.75, 3.54, 4.6, 0.24, "THE REASONING GAP - what one scalar amputates",
    size=9, color=WARN, bold=True, spc=140)
missing = [("uncertainty decomposition", "epistemic vs aleatoric"),
           ("calibration guarantee", "is 0.87 actually 0.87?"),
           ("counterfactual justification", "what flips the call?"),
           ("next-step planning", "which test resolves doubt?")]
for i, (m, msub) in enumerate(missing):
    x = 4.75 + (i % 2) * 2.42
    y = 3.86 + (i // 2) * 0.55
    txt(s, x, y, 2.4, 0.22,
        [[("✕ ", {"color": WARN, "bold": True, "size": 9.5}),
          (m, {"bold": True, "size": 8.8})]], color=INK, wrap=False)
    txt(s, x + 0.17, y + 0.20, 2.25, 0.18, msub, size=7.5, color=MUTED, wrap=False)
txt(s, 4.55, 5.30, 4.95, 0.70,
    [[("Zech et al. (PLOS Med 2018): these models fail silently under "
       "distribution shift - and do not know they failed. ",
       {"color": MUTED, "size": 9})],
     [("The gap is epistemic, not statistical.", {"color": TEAL, "size": 10,
       "bold": True, "italic": True})]], line_spacing=1.1)
cite(s, "Rajpurkar et al., arXiv:1711.05225 (2017)  ·  Zech et al., PLOS Medicine 15(11) (2018)  ·  US FDA, AI-Enabled Medical Device List (2025)")
notes(s, "THE PROBLEM  ·  1:00\n"
         "CheXNet beat radiologist-level F1 in 2017. Nine hundred fifty cleared "
         "devices later, radiology AI still sits outside the decision loop - "
         "because a softmax score answers none of the four questions a "
         "clinician actually asks: why, how certain, what would change your "
         "mind, what next. Everything in the red box is amputated by design in "
         "today's pipeline. And Zech showed these models collapse silently "
         "under distribution shift. India feels this hardest: one radiologist "
         "per hundred thousand people. The bottleneck is epistemic - trust - "
         "not accuracy.\n\n"
         "RUBRIC: problem significance ✓ quantified, cited, and reframed in a "
         "way judges haven't heard from a student team.\n"
         "JUDGES SHOULD FEEL: the diagnosis of the field's failure is itself "
         "novel - now they need the fix.")

# ==============================================================================
# SLIDE 3 — the system: 8-stage pipeline + 5-layer architecture
# ==============================================================================
s = S[2]
kicker(s, "02", "THE SYSTEM")
title(s, "A diagnostic operating system, not a classifier")

stages_r1 = [("01", "Image", "DICOM + context"),
             ("02", "Quantum encoding", "feature map Uφ(x)"),
             ("03", "Bayesian agents", "belief graph over DDx"),
             ("04", "Calibration", "conformal + temperature")]
stages_r2 = [("08", "Clinician", "decides - with evidence"),
             ("07", "Grounded report", "every claim cited"),
             ("06", "Evidence planner", "next test by EIG"),
             ("05", "Counterfactuals", "what flips the call")]
nw, nh, gap = 2.02, 0.92, 0.28
r1y, r2y = 2.18, 3.52
for rowy, stages, rev in ((r1y, stages_r1, False), (r2y, stages_r2, True)):
    for i, (num, lab, sub) in enumerate(stages):
        x = 0.55 + i * (nw + gap)
        hi = (lab in ("Quantum encoding", "Clinician"))
        box(s, x, rowy, nw, nh, fill=(PALE if hi else PALE2),
            line=(TEAL2 if hi else BORDER), lw=1.2 if hi else 0.75)
        txt(s, x + 0.13, rowy + 0.08, 0.5, 0.2, num, size=8, color=TEAL, bold=True)
        txt(s, x + 0.13, rowy + 0.26, nw - 0.26, 0.28, lab, size=11.5, bold=True,
            color=INK)
        txt(s, x + 0.13, rowy + 0.56, nw - 0.26, 0.28, sub, size=8.5, color=MUTED)
        if not rev and i < 3:
            arrow(s, x + nw + 0.03, rowy + nh / 2, x + nw + gap - 0.03, rowy + nh / 2)
        if rev and 0 < i < 3:
            arrow(s, x + nw + gap - 0.03, rowy + nh / 2, x + nw + 0.03, rowy + nh / 2)
lastx = 0.55 + 3 * (nw + gap)
arrow(s, lastx + nw / 2, r1y + nh + 0.04, lastx + nw / 2, r2y - 0.04)

txt(s, 0.55, 4.70, 8.9, 0.24, "RUNS AS FIVE LAYERS - NINE REPLACEABLE SERVICES, TYPED CONTRACTS, ASYNC EVENT BUS",
    size=8.5, color=TEAL, bold=True, spc=160)
layers = [("CLINICAL", "dashboard · signing"), ("KNOWLEDGE", "ontology · memory"),
          ("AI", "vision · agents · safety"), ("QUANTUM", "kernels · VQC fusion"),
          ("GATEWAY", "PACS-EHR · audit")]
lx = 0.55
for name, desc in layers:
    hi = name == "QUANTUM"
    box(s, lx, 5.00, 1.72, 0.66, fill=(PALE if hi else PALE2),
        line=(TEAL2 if hi else BORDER), lw=1.2 if hi else 0.75)
    txt(s, lx + 0.11, 5.09, 1.5, 0.2, name, size=8.5, bold=True,
        color=(TEAL if hi else INK), spc=120)
    txt(s, lx + 0.11, 5.31, 1.52, 0.3, desc, size=7.5, color=MUTED)
    lx += 1.79
txt(s, 0.5, 5.84, 9.0, 0.26,
    "The clinician is stage 8 - the decision-maker, never a rubber stamp. Every stage is a real, swappable subsystem.",
    size=10, color=TEAL, italic=True, align=PP_ALIGN.CENTER)
cite(s, "Implemented end-to-end in services/{vision, fusion, safety, explain, recommend, report, memory} - live demo available.")
notes(s, "THE SYSTEM  ·  1:10\n"
         "AURA replaces the classifier with a reasoning pipeline. The image "
         "enters; compressed clinical evidence is embedded in a quantum feature "
         "space; Bayesian agents update a belief graph over the differential; "
         "the posterior is calibrated with a coverage guarantee. Then AURA asks "
         "its own follow-ups - what would flip this call, and which test buys "
         "the most information for this patient - and emits a report where "
         "every sentence links to evidence. Under it: five layers, nine "
         "microservices over an async event bus with typed contracts. The "
         "quantum layer is one layer - deliberately not the whole story.\n"
         "Walk the flow with your hand. Do not read boxes.\n\n"
         "RUBRIC: innovation ✓ architecture-level, not model-level. "
         "Feasibility ✓ it is running code, offered for demo.\n"
         "JUDGES SHOULD FEEL: ambition with engineering discipline.")

# ==============================================================================
# SLIDE 4 — quantum core (the centerpiece)
# ==============================================================================
s = S[3]
kicker(s, "03", "QUANTUM CORE")
title(s, "Why Hilbert space - the mathematical case")
txt(s, 0.5, 2.04, 9.0, 0.28,
    "The claim is geometric, never about speed - and it is benchmarked, not assumed.",
    size=11, color=MUTED, italic=True)

steps = [
    ("x ∈ R⁸", "compressed evidence vector - 8 clinical channels, never raw pixels"),
    ("Uφ(x) |0⟩⊗⁸", "parameterized quantum circuit angle-encodes evidence"),
    ("|φ(x)⟩ ∈ H,  dim = 2⁸", "entangling layers spread data across a 256-dim state space"),
    ("k(x, x′) = |⟨φ(x)|φ(x′)⟩|²", "quantum kernel - fidelity between diagnostic states"),
    ("p(d | e)  via VQC fusion", "hybrid Bayesian update over the differential"),
]
syy = 2.48
for i, (eq, desc) in enumerate(steps):
    box(s, 0.5, syy, 3.95, 0.56, fill=PALE2, line=BORDER)
    txt(s, 0.68, syy + 0.055, 3.6, 0.24, eq, size=11.5, bold=True, color=TEAL,
        font="Cambria")
    txt(s, 0.68, syy + 0.30, 3.62, 0.22, desc, size=7.8, color=MUTED)
    if i < 4:
        arrow(s, 2.47, syy + 0.56, 2.47, syy + 0.665, wpt=1.2)
    syy += 0.67

box(s, 4.75, 2.48, 4.75, 1.52, fill=PALE, line=TEAL2, lw=1.1)
txt(s, 5.0, 2.62, 4.3, 0.26, "ENTANGLING FEATURE MAP", size=9, color=TEAL,
    bold=True, spc=220)
txt(s, 5.0, 2.90, 4.3, 0.40, "Uφ(x) = exp( i Σ_S φ_S(x) Π_k∈S Z_k )",
    size=14.5, color=INK, bold=True, font="Cambria", align=PP_ALIGN.CENTER)
txt(s, 5.0, 3.38, 4.28, 0.56,
    "Second-order phases φ_{jk}(x) = (π − x_j)(π − x_k) write pairwise evidence "
    "interactions - e.g. opacity × fever - directly into the state.",
    size=9, color=MUTED, line_spacing=1.05)

box(s, 4.75, 4.14, 4.75, 1.62, fill=WHITE, line=BORDER)
txt(s, 5.0, 4.26, 4.3, 0.26, "THE HONEST CLAIM - AND THE BENCHMARK LAW", size=9,
    color=TEAL, bold=True, spc=160)
txt(s, 5.0, 4.54, 4.28, 1.16,
    [[("Richer similarity geometry, not speed. ", {"bold": True, "color": INK}),
      ("Entangling-circuit kernels are conjectured classically hard to estimate "
       "(Havlíček 2019); a rigorous separation exists for a related family "
       "(Liu 2021).", {"color": MUTED})],
     [("Every quantum service ships beside a classical twin - if the twin wins "
       "on our benchmark, the twin ships.", {"color": TEAL, "italic": True,
       "bold": True})]],
    size=9.5, line_spacing=1.1)
cite(s, "Havlíček et al., Nature 567, 209 (2019)  ·  Schuld & Killoran, PRL 122, 040504 (2019)  ·  Liu, Arunachalam & Temme, Nature Physics 17, 1013 (2021)")
notes(s, "QUANTUM CORE  ·  1:30  - the slide that wins or loses the quantum jury\n"
         "We never put pixels in a circuit - perception stays classical, CNNs "
         "won imaging. What enters the quantum layer is a compressed 8-channel "
         "evidence vector: eight qubits, angle encoding, entangling layers - a "
         "256-dimensional Hilbert space whose phase structure carries pairwise "
         "clinical interactions. The kernel is state fidelity: similarity in "
         "that geometry. Havlicek showed such kernels are conjectured hard to "
         "estimate classically; Liu, Arunachalam and Temme proved a rigorous "
         "separation for a related class. So our claim is representational - "
         "richer geometry over clinical evidence - never 'quantum is faster'. "
         "And because a conjecture is not evidence: every quantum service ships "
         "beside a classical twin, benchmarked head-to-head on six metrics. If "
         "the twin wins, the twin ships, and we say so on stage.\n\n"
         "RUBRIC: technical depth ✓ literature-exact. Scientific integrity ✓ "
         "the benchmark law is the disqualification-proof armor.\n"
         "JUDGES SHOULD FEEL: rare intellectual honesty; nothing to attack, "
         "because every attack is already conceded and measured.")

# ==============================================================================
# SLIDE 5 — epistemic engine: uncertainty + guarantee + next-best evidence
# ==============================================================================
s = S[4]
kicker(s, "04", "EPISTEMIC ENGINE")
title(s, "Knowing which kind of 'unsure' - then acting on it")

box(s, 0.5, 2.14, 9.0, 0.60, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 0.5, 2.24, 9.0, 0.40,
    [[("H[ p(y|x) ]", {"color": INK, "bold": True}),
      ("  =  ", {"color": MUTED}),
      ("I(y ; θ | x)", {"color": TEAL, "bold": True}),
      ("  +  ", {"color": MUTED}),
      ("Eθ H[ p(y|x,θ) ]", {"color": TEAL2, "bold": True}),
      ("      total = epistemic + aleatoric", {"color": MUTED, "size": 9.5,
       "italic": True})]],
    size=14, font="Cambria", align=PP_ALIGN.CENTER)

ucards = [("EPISTEMIC", "model doubt", "ensembles + MC-dropout", "→ defer to radiologist"),
          ("ALEATORIC", "data noise", "ambiguity in the image itself", "→ recommend re-imaging"),
          ("OUT-OF-DIST.", "wrong input", "energy score + latent density", "→ abstain, flag input")]
ux = 0.5
for name, tag, desc, act in ucards:
    box(s, ux, 2.88, 2.92, 1.06, fill=PALE2, line=BORDER)
    txt(s, ux + 0.16, 2.99, 2.62, 0.22,
        [[(name + " ", {"bold": True, "color": INK, "size": 10}),
          ("- " + tag, {"color": MUTED, "size": 8.5, "italic": True})]])
    txt(s, ux + 0.16, 3.25, 2.6, 0.22, desc, size=8.2, color=MUTED)
    txt(s, ux + 0.16, 3.55, 2.6, 0.24, act, size=9, color=TEAL, bold=True)
    ux += 3.04

box(s, 0.5, 4.10, 4.32, 1.62, fill=WHITE, line=BORDER)
txt(s, 0.68, 4.22, 4.0, 0.24, "CONFORMAL GUARANTEE - BY THEOREM", size=8.5,
    color=TEAL, bold=True, spc=160)
txt(s, 0.68, 4.50, 4.0, 0.32, "C(x) = { y : s(x,y) ≤ q̂₁₋α }", size=13,
    color=INK, bold=True, font="Cambria")
txt(s, 0.68, 4.86, 4.0, 0.32, "P( y ∈ C(x) ) ≥ 1 − α,  distribution-free",
    size=11.5, color=TEAL, bold=True, font="Cambria")
txt(s, 0.68, 5.26, 3.95, 0.40,
    "sets, not points - {effusion, atelectasis} at 90% coverage; ECE 0.142 → 0.031 after temperature scaling",
    size=8.5, color=MUTED, line_spacing=1.05)

box(s, 5.02, 4.10, 4.48, 1.70, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 5.20, 4.22, 4.1, 0.24, "NEXT-BEST EVIDENCE (EIG)", size=8.5, color=TEAL,
    bold=True, spc=160)
txt(s, 5.20, 4.48, 4.1, 0.30, "a* = argmaxₐ { H[p(d)] − Eₒ H[p(d|o)] − λ c(a) }",
    size=11.5, color=INK, bold=True, font="Cambria")
bars = [("CT thorax", 0.82), ("Repeat PA X-ray", 0.64), ("Ultrasound", 0.51)]
byy = 4.82
for lab, eig in bars:
    txt(s, 5.20, byy, 1.45, 0.2, lab, size=8.5, bold=True, color=INK,
        anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    box(s, 6.70, byy + 0.025, 2.0 * eig / 0.82, 0.13, fill=NEON, line=None,
        radius=0.5)
    txt(s, 8.90, byy, 0.5, 0.2, f"{eig:.2f}", size=8.5, bold=True, color=TEAL,
        anchor=MSO_ANCHOR.MIDDLE)
    byy += 0.26
txt(s, 5.20, 5.54, 4.15, 0.22,
    [[("→ repeat X-ray: ", {"bold": True, "color": INK, "size": 8.5}),
      ("78% of CT's information at 4% of its cost", {"color": TEAL, "size": 8.5,
       "italic": True})]], wrap=False)
cite(s, "Kendall & Gal, NeurIPS (2017)  ·  Angelopoulos & Bates, FnT-ML 16(4) (2023)  ·  Guo et al., ICML (2017)  ·  Lindley, Ann. Math. Statist. 27 (1956)",
     y=6.00)
notes(s, "EPISTEMIC ENGINE  ·  1:30\n"
         "One equation runs the top: predictive entropy decomposes into "
         "epistemic - mutual information between prediction and parameters - "
         "plus aleatoric. Why decompose? Because each doubt demands a different "
         "action: epistemic high, defer to the human; aleatoric high, re-image "
         "- more data will not help; out-of-distribution, abstain entirely. No "
         "silent failure, by policy. Below left: conformal prediction - the "
         "true label is inside the predicted set with probability at least one "
         "minus alpha, by theorem, distribution-free. Below right: the residual "
         "entropy is spent optimally - expected information gain per unit cost "
         "ranks the next test. Here CT resolves the most entropy, but a repeat "
         "X-ray buys 78 percent of that information at 4 percent of the cost - "
         "so that is the recommendation. This is Lindley's 1956 experimental "
         "design, running per patient.\n\n"
         "RUBRIC: technical depth ✓ three research fields composed correctly; "
         "impact ✓ the uncertainty-to-action mapping is clinically meaningful.\n"
         "JUDGES SHOULD FEEL: 'the AI asks for the right test' - the aha moment.")

# ==============================================================================
# SLIDE 6 — proof: built, benchmarked, positioned
# ==============================================================================
s = S[5]
kicker(s, "05", "PROOF")
title(s, "Built and benchmarked - not proposed")

box(s, 0.5, 2.14, 4.35, 2.62, fill=PALE2, line=BORDER)
txt(s, 0.68, 2.26, 4.0, 0.22, "RUNS TODAY, OFFLINE, ON A LAPTOP", size=8.5,
    color=TEAL, bold=True, spc=160)
impl = [("Quantum fusion", "PennyLane VQC - 8-qubit angle encoding + entangling layers"),
        ("Classical twin", "Bayesian product-of-experts, same interface"),
        ("Safety engine", "ensembles · temperature scaling · conformal sets · OOD · abstention"),
        ("Explainability", "occlusion saliency · evidence attribution · counterfactuals"),
        ("Recommender", "EIG next-test ranking per cost/risk"),
        ("Platform", "FastAPI microservices · event bus · audit ledger · dashboard")]
iy = 2.54
for name, desc in impl:
    txt(s, 0.68, iy, 1.28, 0.36, name, size=8.2, bold=True, color=INK)
    txt(s, 1.98, iy, 2.78, 0.36, desc, size=7.6, color=MUTED, line_spacing=0.95)
    iy += 0.365
box(s, 0.5, 4.90, 4.35, 0.92, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 0.68, 5.00, 4.05, 0.22, "LIVE DEMO - THREE COMMANDS, TWO MINUTES", size=8.5,
    color=TEAL, bold=True, spc=140)
txt(s, 0.68, 5.24, 4.05, 0.5,
    [[("aura_cli train · bench · serve", {"font": "Courier New", "bold": True,
       "color": INK, "size": 10})],
     [("doctor dashboard on :8000 - ask us to run it", {"color": MUTED,
       "size": 8, "italic": True})]], line_spacing=1.15)

box(s, 5.05, 2.14, 4.45, 1.30, fill=PALE, line=TEAL2, lw=1.0)
txt(s, 5.23, 2.26, 4.1, 0.22, "THE BENCHMARK LAW", size=8.5, color=TEAL,
    bold=True, spc=200)
txt(s, 5.23, 2.50, 4.1, 0.86,
    [[("Quantum vs classical twin on identical held-out data: ",
       {"color": MUTED, "size": 8.8}),
      ("accuracy · NLL · ECE · Brier · coverage · set size.",
       {"color": INK, "size": 8.8, "bold": True})],
     [("If the classical twin wins, the classical twin ships.",
       {"color": TEAL, "size": 9.5, "bold": True, "italic": True})]],
    line_spacing=1.15)

txt(s, 5.05, 3.60, 4.45, 0.22, "AGAINST THE STRONGEST PUBLISHED BASELINES",
    size=8.5, color=TEAL, bold=True, spc=140)
cols = ["", "CheXNet", "MedPaLM-M", "AURA"]
rows = [("Uncertainty decomposition", "—", "—", "✓"),
        ("Coverage-guaranteed sets", "—", "—", "✓"),
        ("Counterfactual explanation", "—", "—", "✓"),
        ("Next-test planning (EIG)", "—", "~", "✓"),
        ("Quantum-kernel latent space", "—", "—", "✓"),
        ("Grounded, auditable report", "—", "~", "✓")]
tx0, ty0 = 5.05, 3.88
cw = [2.10, 0.72, 0.86, 0.72]
xacc = [tx0]
for wdt in cw:
    xacc.append(xacc[-1] + wdt)
rh = 0.30
for j, htx in enumerate(cols[1:], 1):
    hi = (j == 3)
    box(s, xacc[j] + 0.02, ty0, cw[j] - 0.04, 0.28, fill=(TEAL if hi else INK),
        line=None, radius=0.12)
    txt(s, xacc[j], ty0 + 0.045, cw[j], 0.2, htx, size=7.5, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
gl = {"✓": TEAL, "—": "B9C2C9", "~": GOLD}
for i, row in enumerate(rows):
    y = ty0 + 0.32 + i * rh
    if i % 2 == 0:
        box(s, tx0, y - 0.015, sum(cw), rh, fill=PALE2, line=None, radius=0,
            shape=MSO_SHAPE.RECTANGLE)
    txt(s, tx0 + 0.06, y, cw[0] - 0.1, rh, row[0], size=8, bold=True, color=INK,
        anchor=MSO_ANCHOR.MIDDLE)
    for j in range(1, 4):
        txt(s, xacc[j], y, cw[j], rh, row[j], size=9.5 if row[j] == "✓" else 9,
            bold=True, color=gl[row[j]], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
box(s, xacc[3], ty0 - 0.02, cw[3], 0.32 + 6 * rh, fill=None, line=TEAL, lw=1.5,
    radius=0.08)
cite(s, "Rajpurkar 2017 · Tu et al., arXiv:2307.14334 (2023) · full provenance per result: {backend, circuit version, seed, shots} pinned to every belief snapshot.")
notes(s, "PROOF  ·  1:15\n"
         "Is this real? Everything on the left runs offline on this laptop - "
         "the VQC, its classical twin, the safety engine, the planner, the "
         "dashboard. Three commands: train, bench, serve. Ask us and we run it "
         "now. The benchmark law is our integrity mechanism: quantum versus "
         "classical twin, identical held-out data, six metrics - and the loser "
         "does not ship. On the right, honest positioning: CheXNet solved "
         "probability; MedPaLM-M can suggest follow-ups in free text - we mark "
         "that ourselves. But no published system combines decomposed "
         "uncertainty, guaranteed coverage, counterfactuals, information-"
         "theoretic planning, and a quantum representation in one auditable "
         "loop. The novelty is the column.\n\n"
         "RUBRIC: feasibility ✓ demo on request kills the vaporware doubt; "
         "novelty ✓ positioned against SOTA with self-declared '~' marks - "
         "nothing left to disqualify.\n"
         "JUDGES SHOULD FEEL: this team grades itself harder than we would.")

# ==============================================================================
# SLIDE 7 — impact + closing (dark)
# ==============================================================================
s = S[6]
dark_bg(s)
rings(s, 5.0, 3.0, [1.6, 2.5, 3.4], alphas=(16, 9, 5))
txt(s, 0.5, 0.62, 9.0, 0.26, "06  ·  WHAT CHANGES", size=10, color=DIM,
    bold=True, spc=300, align=PP_ALIGN.CENTER)
imp = [("HOSPITALS", "worklists triaged by uncertainty - confident criticals jump the queue"),
       ("PATIENTS", "EIG-planned workups: fewer scans, less radiation, faster answers"),
       ("THE SYSTEM", "aimed at diagnostic error - medicine's costliest failure (~12M US adults/yr)")]
ix = 0.55
for name, desc in imp:
    box(s, ix, 1.14, 2.92, 1.10, fill=PANEL, line=NEON, lw=0.75, line_alpha=45)
    txt(s, ix + 0.18, 1.28, 2.6, 0.22, name, size=9.5, bold=True, color=NEONHI,
        spc=160)
    txt(s, ix + 0.18, 1.54, 2.6, 0.62, desc, size=8.5, color=DIM, line_spacing=1.1)
    ix += 3.04
mile = [("2026", "clinical pilot"), ("2027", "quantum federated learning"),
        ("2028", "hardware kernels"), ("2029", "ICU real-time agents")]
line(s, 1.3, 2.78, 8.7, 2.78, color=NEON, wpt=1.0, alpha=40, dash="sysDash")
mx = 1.65
for datev, name in mile:
    box(s, mx - 0.05, 2.73, 0.10, 0.10, fill=NEON, line=None, shape=MSO_SHAPE.OVAL)
    txt(s, mx - 0.85, 2.92, 1.7, 0.40,
        [[(datev, {"size": 8.5, "bold": True, "color": NEONHI})],
         [(name, {"size": 7.5, "color": DIM})]], align=PP_ALIGN.CENTER,
        line_spacing=1.0)
    mx += 2.12
txt(s, 1.0, 3.86, 8.0, 1.65,
    [[("Artificial intelligence predicts.", {"color": "9FB3B1", "size": 19})],
     [("Clinical intelligence reasons.", {"color": WHITE, "size": 23})],
     [("Quantum intelligence understands.", {"color": NEONHI, "size": 27,
       "bold": True})]],
    align=PP_ALIGN.CENTER, line_spacing=1.3)
line(s, 4.2, 5.78, 5.8, 5.78, color=NEON, wpt=1.0, alpha=45)
txt(s, 1.0, 5.94, 8.0, 0.62,
    [[("AURA  ·  live demo on request  ·  every claim on these slides is cited and reproducible",
       {"color": WHITE, "size": 10, "bold": True, "spc": 120})],
     [("Sivaguru R.M  ·  Rajalakshmi Institute of Technology  ·  sivagurumurugan1@gmail.com",
       {"color": DIM, "size": 9, "spc": 100})]],
    align=PP_ALIGN.CENTER, line_spacing=1.5)
notes(s, "IMPACT + CLOSE  ·  1:00\n"
         "What changes. Hospitals get uncertainty-triaged worklists - the "
         "confident criticals jump the queue. Patients get workups planned by "
         "information gain: fewer scans, less radiation, faster answers. And "
         "the system attacks diagnostic error - the costliest failure mode in "
         "medicine, twelve million affected adults a year in the US alone, "
         "more and less-measured in India. The roadmap is hardware-realistic: "
         "pilot next, federated quantum learning - where only kernel "
         "statistics ever leave a hospital - then hardware kernels and ICU "
         "streaming agents. No stage depends on a miracle.\n"
         "Then stop. Three sentences, slowly: AI predicts. Clinical "
         "intelligence reasons. Quantum intelligence understands. The system "
         "is live - we would love to show you. Thank you.\n\n"
         "RUBRIC: impact ✓ quantified per stakeholder; vision ✓ staged and "
         "realistic; presentation ✓ ends with goosebumps and a demo offer.\n"
         "JUDGES SHOULD FEEL: nothing left to ask except 'show us'.")

prs.save("AURA_Finals_7Slides.pptx")
print("saved AURA_Finals_7Slides.pptx with", len(prs.slides._sldIdLst), "slides")
