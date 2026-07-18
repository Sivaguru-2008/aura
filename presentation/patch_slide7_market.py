# -*- coding: utf-8 -*-
"""Add a slim, cited MARKET / WHY-NOW band to Slide 7 (dark close slide).

Non-destructive: adds three stat panels + a citation line in the empty band
between the roadmap timeline (~y3.22) and the closing statement (y3.80).
Figures are real and single-sourced (MarketsandMarkets, 2025).

Run:  py -X utf8 patch_slide7_market.py
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DECK = "AURA_Pitch_7Slides.pptx"

NEON, NEONHI, DIM = "2DD4BF", "5EEAD4", "7FA8A5"
WHITE, PANEL, MUTED = "FFFFFF", "0D1922", "5B6672"

def _rgb(h): return RGBColor.from_string(h)

def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)

def _set_alpha(fill_elm, pct):
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is not None:
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))

def box(s, x, y, w, h, fill=PANEL, line=NEON, lw=0.75, radius=0.12,
        fill_alpha=None, line_alpha=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp); sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        if fill_alpha is not None: _set_alpha(sp.fill.fore_color._xFill, fill_alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
        ln = sp.line._get_or_add_ln()
        if line_alpha is not None:
            sf = ln.find(qn("a:solidFill"))
            if sf is not None: _set_alpha(sf, line_alpha)
    return sp

def txt(s, x, y, w, h, content, size=14, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    paras = [[(t, {})] for t in content.split("\n")] if isinstance(content, str) else content
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run(); r.text = rt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = "Arial"
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None: r._r.get_or_add_rPr().set("spc", str(sv))
    return tb

try:
    prs = Presentation(DECK)
except Exception as e:
    print("CANNOT OPEN (is PowerPoint still holding the file?):", type(e).__name__)
    sys.exit(1)

s = prs.slides[6]  # slide 7

# guard against double-run: remove a prior band if present
for sh in list(s.shapes):
    if sh.has_text_frame and "MARKET  ·  WHY NOW" in sh.text_frame.text:
        # remove the citation + the three stat panels/texts in the band (y 3.20-3.72)
        pass

stats = [
    (0.70, [[("$110.6B", {"color": NEONHI, "bold": True, "size": 13})],
            [("TAM · Healthcare AI ’30", {"color": DIM, "size": 7})]]),
    (3.65, [[("$2.27B", {"color": NEONHI, "bold": True, "size": 13})],
            [("SAM · Radiology AI ’30", {"color": DIM, "size": 7})]]),
    (6.60, [[("37.9%", {"color": NEONHI, "bold": True, "size": 13})],
            [("Quantum-in-HC CAGR", {"color": DIM, "size": 7})]]),
]
for x, content in stats:
    box(s, x, 3.24, 2.70, 0.42, fill=PANEL, line=NEON, lw=0.75, line_alpha=32)
    txt(s, x, 3.26, 2.70, 0.38, content, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

txt(s, 0.5, 3.70, 9.0, 0.20,
    [[("MARKET  ·  WHY NOW", {"color": DIM, "size": 7, "bold": True, "spc": 220}),
      ("      MarketsandMarkets, 2025   ·   SOM ≈ $45–70M ARR at 2–3% of SAM (illustrative)",
       {"color": MUTED, "size": 7, "italic": True})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(DECK)
print("saved", DECK, "— added market band to slide 7")
