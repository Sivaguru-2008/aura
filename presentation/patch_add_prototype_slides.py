# -*- coding: utf-8 -*-
"""Insert two 'working prototype' slides after slide 4 of AURA_Template_Final.pptx.

Each new slide duplicates slide 4's chrome (header logos + footer placeholder),
clears the mockup body, and shows a REAL high-res screenshot of the running
console plus a deep, numbered walkthrough of every subsystem:

  new slide 5  — CONFIDENT READ: scan -> vision -> quantum fusion -> report
  new slide 6  — CALIBRATED DOUBT: the abstention / safety story

Screenshots come from ../media/proto_console.png and ../media/proto_abstain.png
(captured headless from presentation/AURA_Prototype.html).

Run:  py -X utf8 patch_add_prototype_slides.py
Idempotent-ish: re-running appends again, so start from the backup if needed.
"""
from __future__ import annotations
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DECK = "AURA_Template_Final.pptx"
IMG_CONSOLE = "../media/proto_console.png"
IMG_ABSTAIN = "../media/proto_abstain.png"

# deck design tokens (from build_pitch.py)
INK, MUTED = "0F172A", "5B6672"
TEAL, TEAL2, NEON, NEONHI = "0E7490", "0D9488", "2DD4BF", "5EEAD4"
PALE, PALE2, BORDER = "ECF8F7", "F6FAFA", "C6E4E1"
WHITE, DARKBG, DIM = "FFFFFF", "05090E", "7FA8A5"
WARN, GOLD = "C05B52", "B98A2F"


def _rgb(h):
    return RGBColor.from_string(h)


def _strip_style(sp):
    st = sp._element.find(qn("p:style"))
    if st is not None:
        sp._element.remove(st)


def box(s, x, y, w, h, fill=PALE2, line=BORDER, lw=0.75, radius=0.10,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _strip_style(sp)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(lw)
    return sp


def txt(s, x, y, w, h, content, size=14, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial", spc=None,
        wrap=True, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(content, str):
        paras = [[(t, {})] for t in content.split("\n")]
    else:
        paras = [[(it, {})] if isinstance(it, str) else it for it in content]
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for rt, opts in runs:
            r = p.add_run()
            r.text = rt
            f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.name = opts.get("font", font)
            f.color.rgb = _rgb(opts.get("color", color))
            sv = opts.get("spc", spc)
            if sv is not None:
                r._r.get_or_add_rPr().set("spc", str(sv))
    return tb


def marker(s, cx, cy, n, d=0.28):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                            Inches(d), Inches(d))
    _strip_style(sp)
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(NEON)
    sp.line.color.rgb = _rgb(DARKBG)
    sp.line.width = Pt(1.0)
    txt(s, cx - d / 2, cy - d / 2, d, d, str(n), size=11, color=DARKBG, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def dup_slide(prs, index):
    """Duplicate slide at `index`, preserving pictures/logos via rel remap."""
    src = prs.slides[index]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):            # drop layout-copied placeholders
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:                  # deep-copy every shape element
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    id_map = {}                            # copy rels (images etc.), remap rIds
    for rId, rel in src.part.rels.items():
        rt = rel.reltype
        if rt.endswith("slideLayout") or rt.endswith("notesSlide"):
            continue
        if rel.is_external:
            id_map[rId] = new.part.relate_to(rel._target, rt, is_external=True)
        else:
            id_map[rId] = new.part.relate_to(rel.target_part, rt)
    for el in new.shapes._spTree.iter():
        for a in (qn("r:embed"), qn("r:link")):
            v = el.get(a)
            if v in id_map:
                el.set(a, id_map[v])
    return new


def clear_body(s):
    """Remove everything except the header logos (top<1.5in) and the footer."""
    for sh in list(s.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        top = (sh.top or 0) / 914400.0
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and top < 1.5:
            continue
        sh._element.getparent().remove(sh._element)


def build_hero(s, kick, ttl, img, cap, markers, heading, steps, note):
    clear_body(s)
    # kicker + title (match slide-4 positions)
    txt(s, 0.4, 1.22, 9.2, 0.26,
        [[("03", {"color": MUTED}), ("   .   ", {"color": MUTED}), (kick, {})]],
        size=11, color=TEAL, bold=True, spc=200)
    txt(s, 0.4, 1.52, 9.2, 0.52, ttl, size=23, color=INK, bold=True, spc=-10)
    # hero screenshot 1600x1000 -> 1.6:1
    PX, PY, PW = 0.4, 2.14, 5.75
    PH = PW * 1000.0 / 1600.0            # 3.594
    pic = s.shapes.add_picture(img, Inches(PX), Inches(PY), width=Inches(PW))
    pic.line.color.rgb = _rgb(TEAL2)
    pic.line.width = Pt(1.5)
    txt(s, PX, PY + PH + 0.07, PW, 0.24,
        [[("localhost:8000 . AURA console", {"color": TEAL, "size": 8, "bold": True}),
          ("   " + cap, {"color": MUTED, "size": 8, "italic": True})]])
    for n, fx, fy in markers:
        marker(s, PX + fx * PW, PY + fy * PH, n)
    # right-column deep walkthrough — on a card so the deck's background motif
    # sits behind a fill (matches slide 3) instead of running through the text
    CX, CW = 6.42, 3.16
    box(s, CX - 0.2, 1.98, CW + 0.38, 4.4, fill=PALE2, line=BORDER, lw=0.75, radius=0.04)
    txt(s, CX, 2.14, CW, 0.24, heading, size=9.5, color=TEAL, bold=True, spc=120)
    sy = 2.52
    for n, head, body in steps:
        txt(s, CX, sy, CW, 0.56,
            [[("%d  " % n, {"color": TEAL, "bold": True, "size": 10}),
              (head, {"color": INK, "bold": True, "size": 8.6, "spc": 30})],
             [(body, {"color": MUTED, "size": 8})]],
            line_spacing=1.04)
        sy += 0.615
    s.notes_slide.notes_text_frame.text = note


def remove_slide(prs, slide):
    """Delete a slide (drop its sldId + relationship)."""
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rId = sid.get(qn("r:id"))
        if prs.part.rels[rId].target_part is slide.part:
            prs.part.drop_rel(rId)
            lst.remove(sid)
            return True
    return False


def copy_bg(src, dst):
    """Copy the deck's background motif (a blipFill image in <p:bg>) from src
    slide to dst slide, remapping the image relationship."""
    scs = src._element.find(qn("p:cSld"))
    dcs = dst._element.find(qn("p:cSld"))
    sbg = scs.find(qn("p:bg"))
    if sbg is None:
        return
    old = dcs.find(qn("p:bg"))
    if old is not None:
        dcs.remove(old)
    nbg = copy.deepcopy(sbg)
    for blip in nbg.iter(qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        if rid and rid in src.part.rels:
            rel = src.part.rels[rid]
            blip.set(qn("r:embed"), dst.part.relate_to(rel.target_part, rel.reltype))
    dcs.insert(0, nbg)                      # <p:bg> must be first child of <p:cSld>


def swap_slide4(prs):
    """Replace slide 4's hand-drawn mock dashboard (right side) with the real
    console screenshot. Keeps the '6 STEPS' walkthrough panel on the left."""
    s = prs.slides[3]
    for sh in list(s.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        left = (sh.left or 0) / 914400.0
        top = (sh.top or 0) / 914400.0
        if left >= 5.0 and 2.4 <= top <= 6.1:      # the mock dashboard region
            sh._element.getparent().remove(sh._element)
    # "14 cases" in the subtitle -> match the screenshot's live worklist
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "14 cases" in r.text:
                        r.text = r.text.replace("14 cases", "live worklist")
    PX, PY, PW = 5.05, 2.62, 4.55
    PH = PW * 1000.0 / 1600.0
    pic = s.shapes.add_picture(IMG_CONSOLE, Inches(PX), Inches(PY), width=Inches(PW))
    pic.line.color.rgb = _rgb(TEAL2)
    pic.line.width = Pt(1.5)
    txt(s, PX, PY + PH + 0.09, PW, 0.24,
        [[("real console . localhost:8000", {"color": TEAL, "size": 8, "bold": True}),
          ("  - not a mockup: every value is live model output",
           {"color": MUTED, "size": 8, "italic": True})]])


prs = Presentation(DECK)

# ---- SLIDE 5: the confident read / full pipeline --------------------------
s5 = dup_slide(prs, 3)
build_hero(
    s5,
    "THE WORKING PROTOTYPE - A REAL, RUNNING CONSOLE",
    "Inside the console: scan to calibrated sign-off",
    IMG_CONSOLE,
    "- every value on screen is live model output, offline & CPU-only",
    [(1, 0.075, 0.40), (2, 0.30, 0.34), (3, 0.30, 0.83),
     (4, 0.60, 0.42), (5, 0.88, 0.24), (6, 0.88, 0.83)],
    "HOW IT WORKS - END TO END",
    [
        (1, "WORKLIST - TRIAGED BY UNCERTAINTY",
         "14 cases ranked by risk; low-confidence studies flip to ABSTAINED instead of surfacing a guess."),
        (2, "VISION + SALIENCY",
         "A CNN localises the finding on the CXR and saliency shows where it looked - here pneumothorax at 0.99."),
        (3, "QUANTUM FUSION -> CALIBRATED",
         "An 8-qubit VQC (512 shots) fuses the evidence; the posterior is 13.8x better calibrated than the classical twin (ECE 0.020 vs 0.276)."),
        (4, "EVIDENCE -> REASONING",
         "Eight channels converge into one belief; hover any node for a live counterfactual - remove it and the top belief drops by 78%."),
        (5, "SAFETY ENVELOPE",
         "Epistemic, aleatoric and OOD-energy dials plus a 90% conformal set carrying a coverage guarantee."),
        (6, "GROUNDED REPORT & SIGN-OFF",
         "The report writes itself and every sentence traces to its evidence; the doctor accepts, edits or signs - authority stays human."),
    ],
    "WORKING PROTOTYPE - THE CONFIDENT READ  .  1:00\n"
    "This is the actual running console, not a mockup - every number is live "
    "model output, offline and CPU-only. Walk the pipeline: the worklist is "
    "triaged by uncertainty; the vision model localises the finding with "
    "saliency; the 8-qubit quantum fusion produces a calibrated differential - "
    "99.5% pneumothorax, and calibrated means it. The evidence-to-reasoning "
    "graph exposes each finding's contribution with live counterfactuals on "
    "hover. The safety panel reports epistemic, aleatoric and OOD energy with a "
    "90% conformal set. Finally the grounded report writes itself, every "
    "sentence traced to evidence, and the doctor signs.\n\n"
    "JUDGES SHOULD FEEL: this is real, end-to-end, running software.")

# ---- SLIDE 6: calibrated doubt / abstention -------------------------------
s6 = dup_slide(prs, 3)
build_hero(
    s6,
    "CALIBRATED DOUBT - WHEN AURA REFUSES TO GUESS",
    "The product is the license to say “I don’t know”",
    IMG_ABSTAIN,
    "- same console, an uncertain case: AURA abstains and escalates",
    [(1, 0.30, 0.83), (2, 0.82, 0.30), (3, 0.91, 0.22),
     (4, 0.82, 0.40), (5, 0.60, 0.80), (6, 0.90, 0.83)],
    "WHY THIS IS THE DIFFERENTIATOR",
    [
        (1, "SPREAD DIFFERENTIAL",
         "No diagnosis dominates: COPD 48%, pneumonia 29%, heart failure 15% - the features genuinely overlap."),
        (2, "LARGE CONFORMAL SET (90%)",
         "Three diagnoses stay statistically plausible at the 90% level, so the guaranteed-coverage set is large."),
        (3, "UNCERTAINTY DECOMPOSED",
         "Epistemic 0.138, aleatoric 1.240, OOD 0.86 - high, and reported honestly rather than hidden."),
        (4, "ABSTENTION - NO SILENT GUESS",
         "AURA declines to commit and escalates carrying its full uncertainty state. This is the whole point."),
        (5, "NEXT BEST EVIDENCE (EIG)",
         "It ranks the single most discriminating missing fact - prior films, +0.50 bits at zero cost and risk."),
        (6, "HUMAN-IN-THE-LOOP",
         "The case is handed to the radiologist with the reasoning attached - authority always stays human."),
    ],
    "CALIBRATED DOUBT - THE ABSTENTION  .  0:55\n"
    "Same console, a harder case. Nothing dominates the differential - COPD, "
    "pneumonia and heart failure all remain plausible - so the 90% conformal "
    "set is large and the uncertainty is high. Instead of guessing, AURA "
    "ABSTAINS: it declines to commit and escalates to the radiologist carrying "
    "its full uncertainty state - no silent failure. Then it stays useful: the "
    "planner ranks the next best test by information gain per cost and risk - "
    "prior films, half a bit, free. Authority always stays human. This "
    "abstention is the product; it is what no deployed tool does today.\n\n"
    "JUDGES SHOULD FEEL: this is the moat - it knows what it doesn't know.")

# ---- match the deck: give the new slides the same background motif ----------
copy_bg(prs.slides[3], s5)
copy_bg(prs.slides[3], s6)

# ---- swap slide 4's mock for a real screenshot (after dups copied the mock) -
swap_slide4(prs)

# ---- reorder: place the two new slides right after slide 4 ----------------
lst = prs.slides._sldIdLst
ids = list(lst)                 # [s1..s7, s5new, s6new]
n5, n6 = ids[-2], ids[-1]
lst.remove(n5)
lst.remove(n6)
lst.insert(4, n5)               # 5th slide
lst.insert(5, n6)               # 6th slide

# ---- final edit: drop the old Solution slide (3rd) and the redundant
# confident deep-dive (5th, = s5), leaving a tight 7-slide deck --------------
solution = prs.slides[2]        # "AURA: The Clinical Reasoning Engine"
remove_slide(prs, s5)           # confident deep-dive (overview lives on slide 4)
remove_slide(prs, solution)

prs.save(DECK)
print("saved", DECK, "with", len(prs.slides._sldIdLst), "slides")
